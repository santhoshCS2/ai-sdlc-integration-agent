import os
import json
import re
from typing import Dict, List, Any
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, black, white
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
import logging
from xml.sax.saxutils import escape
from datetime import datetime

# Universal file readers
try:
    from docx import Document
except ImportError:
    Document = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pandas as pd
except ImportError:
    pd = None

from app.agents.architecture.services.github_architecture_service import SystemArchitecture
from app.agents.architecture.services.diagram_generator import ArchitectureDiagramGenerator
from app.agents.architecture.services.layered_diagram_generator import LayeredDataFlowGenerator

logger = logging.getLogger(__name__)

class MockRepoAnalysis:
    """Mock repository analysis object for diagram generation"""
    def __init__(self, repo_data):
        self.tech_stack = {
            'frontend': repo_data.get('frontend_tech', []),
            'backend': repo_data.get('backend_tech', []),
            'languages': list(repo_data.get('languages', {}).keys()),
            'database': repo_data.get('database_tech', [])
        }
        self.api_endpoints = repo_data.get('api_endpoints', [])
        self.components = [f"Component_{i}" for i in range(repo_data.get('components_total', 0))]
        self.backend_structure = {'services': []}
        self.business_logic = []
        self.dependencies = {'production': []}

class GitHubPDFService:
    def __init__(self, output_dir: str = "generated_pdfs"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

        
        self.colors = {
            'primary': HexColor('#2E86AB'),
            'secondary': HexColor('#A23B72'),
            'text': HexColor('#2D3748'),
            'light_gray': HexColor('#F7FAFC'),
            'medium_gray': HexColor('#E2E8F0'),
        }
        
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _sanitize_text(self, text: Any) -> str:
        return "" if not text else escape(str(text))

    def _setup_custom_styles(self):
        self.styles.add(ParagraphStyle(name='CustomTitle', parent=self.styles['Title'], fontSize=24, spaceAfter=30, textColor=self.colors['primary'], alignment=TA_CENTER, fontName='Helvetica-Bold'))
        self.styles.add(ParagraphStyle(name='CustomHeading1', parent=self.styles['Heading1'], fontSize=18, spaceAfter=12, spaceBefore=20, textColor=self.colors['primary'], fontName='Helvetica-Bold'))
        self.styles.add(ParagraphStyle(name='CustomHeading2', parent=self.styles['Heading2'], fontSize=14, spaceAfter=10, spaceBefore=15, textColor=self.colors['secondary'], fontName='Helvetica-Bold'))
        self.styles.add(ParagraphStyle(name='CustomBody', parent=self.styles['Normal'], fontSize=10, spaceAfter=6, textColor=self.colors['text'], alignment=TA_JUSTIFY))
        self.styles.add(ParagraphStyle(name='CustomBullet', parent=self.styles['Normal'], fontSize=10, spaceAfter=4, leftIndent=20, textColor=self.colors['text']))
        self.styles.add(ParagraphStyle(name='GitHubCode', parent=self.styles['Normal'], fontSize=9, fontName='Courier', textColor=self.colors['text'], backColor=self.colors['light_gray'], leftIndent=10, rightIndent=10, spaceAfter=6))
    
    def extract_text_from_file(self, file_path: str) -> str:
        """Universal document reader - supports PDF, DOCX, PPTX, XLSX, TXT, and more"""
        if not os.path.exists(file_path):
            raise ValueError(f"File not found: {file_path}")
        
        ext = os.path.splitext(file_path.lower())[1]
        text = ""
        
        try:
            if ext == '.pdf':
                text = self._extract_pdf_text(file_path)
            elif ext in ['.docx', '.doc']:
                text = self._extract_word_text(file_path)
            elif ext in ['.pptx', '.ppt']:
                text = self._extract_powerpoint_text(file_path)
            elif ext in ['.xlsx', '.xls']:
                text = self._extract_excel_text(file_path)
            elif ext in ['.txt', '.md', '.rtf']:
                text = self._extract_plain_text(file_path)
            else:
                # Try as plain text fallback
                text = self._extract_plain_text(file_path)
            
            # Universal text cleaning
            return self._clean_extracted_text(text)
            
        except Exception as e:
            logger.error(f"Document extraction failed for {file_path}: {str(e)}")
            raise ValueError(f"Failed to read {ext} file: {str(e)}. Ensure required libraries are installed.")
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF files with multiple fallback methods"""
        text = ""
        
        # Try pdfplumber first (better text extraction)
        if pdfplumber:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                if text.strip():
                    return text
            except Exception as e:
                logger.warning(f"pdfplumber failed: {e}")
        
        # Fallback to PyPDF2
        if PyPDF2:
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                if text.strip():
                    return text
            except Exception as e:
                logger.warning(f"PyPDF2 failed: {e}")
        
        raise ImportError("No PDF libraries available. Install: pip install pdfplumber PyPDF2")
    
    def _extract_word_text(self, file_path: str) -> str:
        """Extract text from Word documents with enhanced cleaning"""
        if not Document:
            raise ImportError("python-docx not available. Install: pip install python-docx")
        
        text = ""
        doc = Document(file_path)
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text.strip():
                    text += row_text + "\n"
        
        return text
    
    def _extract_powerpoint_text(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        if not Presentation:
            raise ImportError("python-pptx not available. Install: pip install python-pptx")
        
        text = ""
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        
        return text
    
    def _extract_excel_text(self, file_path: str) -> str:
        """Extract text from Excel files"""
        text = ""
        
        # Try pandas first
        if pd:
            try:
                df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
                for sheet_name, sheet_df in df.items():
                    text += f"\n--- Sheet: {sheet_name} ---\n"
                    text += sheet_df.to_string(index=False) + "\n"
                return text
            except Exception as e:
                logger.warning(f"pandas Excel read failed: {e}")
        
        # Fallback to openpyxl
        if openpyxl:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, data_only=True)
                for sheet_name in wb.sheetnames:
                    text += f"\n--- Sheet: {sheet_name} ---\n"
                    ws = wb[sheet_name]
                    for row in ws.iter_rows(values_only=True):
                        row_text = ' | '.join([str(cell) for cell in row if cell is not None])
                        if row_text.strip():
                            text += row_text + "\n"
                return text
            except Exception as e:
                logger.warning(f"openpyxl failed: {e}")
        
        raise ImportError("No Excel libraries available. Install: pip install pandas openpyxl")
    
    def _extract_plain_text(self, file_path: str) -> str:
        """Extract text from plain text files with encoding detection"""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'ascii']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        # Last resort: read as binary and decode with errors='ignore'
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='ignore')
    
    def _clean_extracted_text(self, text: str) -> str:
        """Universal text cleaning to prevent PDF generation errors"""
        if not text:
            return ""
        
        # Remove PDF artifacts
        if text.startswith('%PDF'):
            text = re.sub(r'%PDF-.*?(\n\n|\n[A-Z])', '', text, flags=re.DOTALL)
        
        # Remove binary/control characters that cause paraparser errors
        text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]', '', text)
        
        # Replace problematic Unicode and special characters
        replacements = {
            '\u2018': "'",  # Left single quote
            '\u2019': "'",  # Right single quote
            '\u201c': '"',  # Left double quote
            '\u201d': '"',  # Right double quote
            '\u2013': '-',  # En dash
            '\u2014': '--', # Em dash
            '\u2026': '...', # Ellipsis
            '\u00a0': ' ',  # Non-breaking space
            '&quot;': '"',  # HTML encoded quote
            '&amp;': '&',   # HTML encoded ampersand
            '&lt;': '<',    # HTML encoded less than
            '&gt;': '>',    # HTML encoded greater than
            '&nbsp;': ' ',  # HTML non-breaking space
            '\u00b0': ' degrees',  # Degree symbol
            '\u2022': '•',  # Bullet point
            '\u00ae': '(R)', # Registered trademark
            '\u00a9': '(C)', # Copyright
            '\u2122': '(TM)', # Trademark
        }
        
        for old, new in replacements.items():
            text = text.replace(old, new)
        
        # Remove HTML tags that might cause issues
        text = re.sub(r'<[^>]+>', '', text)
        
        # Clean up whitespace and line breaks
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        clean_text = '\n'.join(lines)
        
        # Remove any remaining non-ASCII characters that could cause issues
        clean_text = re.sub(r'[^\x20-\x7E\n\r\t•]', ' ', clean_text)
        
        # Clean up multiple spaces
        clean_text = re.sub(r'\s+', ' ', clean_text)
        clean_text = re.sub(r'\n\s*\n', '\n', clean_text)
        
        return clean_text.strip()

    def _setup_custom_styles(self):
        self.styles.add(ParagraphStyle(name='CustomTitle', parent=self.styles['Title'], fontSize=24, spaceAfter=30, textColor=self.colors['primary'], alignment=TA_CENTER, fontName='Helvetica-Bold'))
        self.styles.add(ParagraphStyle(name='CustomHeading1', parent=self.styles['Heading1'], fontSize=18, spaceAfter=12, spaceBefore=20, textColor=self.colors['primary'], fontName='Helvetica-Bold'))
        self.styles.add(ParagraphStyle(name='CustomHeading2', parent=self.styles['Heading2'], fontSize=14, spaceAfter=10, spaceBefore=15, textColor=self.colors['secondary'], fontName='Helvetica-Bold'))
        self.styles.add(ParagraphStyle(name='CustomBody', parent=self.styles['Normal'], fontSize=10, spaceAfter=6, textColor=self.colors['text'], alignment=TA_JUSTIFY))
        self.styles.add(ParagraphStyle(name='CustomBullet', parent=self.styles['Normal'], fontSize=10, spaceAfter=4, leftIndent=20, textColor=self.colors['text']))
        self.styles.add(ParagraphStyle(name='GitHubCode', parent=self.styles['Normal'], fontSize=9, fontName='Courier', textColor=self.colors['text'], backColor=self.colors['light_gray'], leftIndent=10, rightIndent=10, spaceAfter=6))
        self.styles.add(ParagraphStyle(name='PlainASCII', parent=self.styles['Normal'], fontSize=9, fontName='Courier-Bold', textColor=black, leftIndent=0, rightIndent=0, spaceAfter=2, spaceBefore=0, alignment=TA_LEFT))

    def analyze_repo_from_object(self, repo_analysis) -> Dict:
        """Convert repo_analysis object to expected format with real data extraction"""
        try:
            logger.info(f"🔍 Analyzing RepositoryAnalysis object: {type(repo_analysis)}")
            
            # Extract data from RepositoryAnalysis dataclass
            tech_stack = getattr(repo_analysis, 'tech_stack', {})
            languages = tech_stack.get('languages', [])
            frontend_tech = tech_stack.get('frontend', [])
            backend_tech = tech_stack.get('backend', [])
            database_tech = tech_stack.get('database', [])
            
            # Get folder structure and calculate file count
            folder_structure = getattr(repo_analysis, 'folder_structure', {})
            actual_file_count = 0
            for folder, info in folder_structure.items():
                if isinstance(info, dict) and 'files' in info:
                    actual_file_count += len(info['files'])
            
            # Get components count
            components = getattr(repo_analysis, 'components', [])
            actual_components = len(components)
            
            # Get API endpoints
            api_endpoints = getattr(repo_analysis, 'api_endpoints', [])
            
            # Get build tools
            build_tools = getattr(repo_analysis, 'build_tools', [])
            
            logger.info(f"📊 Extracted data - Files: {actual_file_count}, Components: {actual_components}, APIs: {len(api_endpoints)}")
            logger.info(f"📊 Tech stack - Languages: {languages}, Frontend: {frontend_tech}, Backend: {backend_tech}")
            
            # Convert languages list to dict format expected by PDF service
            languages_dict = {}
            for lang in languages:
                languages_dict[lang] = 100 // max(len(languages), 1)  # Equal distribution
            
            # Infer missing technologies with intelligent defaults
            if not frontend_tech and any(lang in ['JavaScript', 'TypeScript'] for lang in languages):
                frontend_tech = ['React', 'Vite']
            
            if not backend_tech:
                if any(lang in ['Python'] for lang in languages):
                    backend_tech = ['Python', 'FastAPI', 'Pydantic']
                elif api_endpoints or frontend_tech:  # If we have frontend or APIs, assume Python backend
                    backend_tech = ['Python', 'FastAPI', 'Uvicorn']
                    languages_dict['Python'] = 60  # Add Python to languages
            
            if not database_tech and (backend_tech or api_endpoints):
                database_tech = ['PostgreSQL', 'SQLAlchemy']
            
            result = {
                'languages': languages_dict,
                'frontend_tech': frontend_tech,
                'backend_tech': backend_tech,
                'database_tech': database_tech,
                'build_tools': build_tools,
                'components_total': actual_components,
                'pages_total': 0,  # Not tracked in RepositoryAnalysis
                'api_endpoints': [{
                    'method': ep.method,
                    'path': ep.path,
                    'purpose': ep.purpose,
                    'file_location': ep.file_location
                } for ep in api_endpoints],
                'patterns': self._infer_architecture_pattern(frontend_tech, backend_tech),
                'file_count': actual_file_count,
                'total_lines': 0,  # Not tracked in RepositoryAnalysis
                'folder_structure': folder_structure
            }
            
            logger.info(f"🎯 Final repo analysis result: file_count={result['file_count']}, languages={len(result['languages'])}, components={result['components_total']}, apis={len(result['api_endpoints'])}")
            return result
        except Exception as e:
            logger.error(f"❌ Error analyzing repo object: {str(e)}")
            # Log the repo_analysis object structure for debugging
            if hasattr(repo_analysis, '__dict__'):
                logger.error(f"📋 Repo analysis attributes: {list(repo_analysis.__dict__.keys())}")
            return self._get_default_analysis()
    
    def _extract_entities_from_prd(self, prd_content: str) -> List[str]:
        """Extract database entities from PRD schema section"""
        if not prd_content:
            return []
        
        entities = set()
        
        # Extract from database schema table definitions
        schema_patterns = [
            r'(?:Table|Entity)\s+([A-Za-z]+)\s+(?:Fields|Columns)',
            r'([A-Za-z]+)\s+(?:table|entity)\s*:?\s*(?:Fields|Columns)',
            r'\b([A-Z][a-z]+)\s+(?:user_id|hotel_id|room_id|booking_id)',  # Table with ID fields
        ]
        
        for pattern in schema_patterns:
            matches = re.findall(pattern, prd_content, re.IGNORECASE)
            entities.update([m.lower() for m in matches])
        
        # Extract from explicit API endpoints
        api_patterns = [
            r'/(?:api/)?([a-zA-Z]+)(?:/|\s|$)',  # /api/users, /hotels
            r'/auth/([a-zA-Z]+)',  # /auth/register, /auth/login
        ]
        
        for pattern in api_patterns:
            matches = re.findall(pattern, prd_content)
            for match in matches:
                if match.lower() not in ['api', 'auth', 'register', 'login']:
                    entities.add(match.lower())
        
        # Clean and validate entities
        valid_entities = []
        for entity in entities:
            if len(entity) >= 3 and entity.isalpha():
                valid_entities.append(entity)
        
        return valid_entities[:6]
    
    def _infer_backend_technologies(self) -> List[str]:
        """Dynamically infer backend technologies based on repository analysis"""
        backend_tech = []
        
        # Check languages to infer backend frameworks
        languages = self._repo_analysis.get('languages', {})
        
        if 'Python' in languages:
            backend_tech.extend(['Python', 'FastAPI'])
        elif 'JavaScript' in languages or 'TypeScript' in languages:
            backend_tech.extend(['Node.js', 'Express'])
        elif 'Java' in languages:
            backend_tech.extend(['Java', 'Spring Boot'])
        elif 'Go' in languages:
            backend_tech.extend(['Go', 'Gin'])
        elif 'C#' in languages:
            backend_tech.extend(['C#', '.NET'])
        else:
            # Default to Python if we have frontend or APIs but no clear backend language
            if self._repo_analysis.get('frontend_tech') or self._prd_analysis.get('api_methods'):
                backend_tech.extend(['Python', 'FastAPI'])
        
        # Add database if we have backend
        if backend_tech:
            backend_tech.append('PostgreSQL')
        
        return backend_tech
    
    def _infer_programming_languages(self) -> List[str]:
        """Dynamically infer programming languages based on repository analysis"""
        languages = []
        
        # Check if we have frontend components
        if self._repo_analysis.get('frontend_tech') or self._repo_analysis.get('components_total', 0) > 0:
            languages.append('JavaScript')
        
        # Check if we have backend or APIs
        if self._repo_analysis.get('backend_tech') or len(self._prd_analysis.get('api_methods', [])) > 0:
            languages.append('Python')
        
        # Check folder structure for file extensions
        folder_structure = self._repo_analysis.get('folder_structure', {})
        for folder, info in folder_structure.items():
            if isinstance(info, dict) and 'files' in info:
                for file in info['files']:
                    if file.endswith(('.py',)) and 'Python' not in languages:
                        languages.append('Python')
                    elif file.endswith(('.js', '.jsx', '.ts', '.tsx')) and 'JavaScript' not in languages:
                        languages.append('JavaScript')
                    elif file.endswith(('.java',)) and 'Java' not in languages:
                        languages.append('Java')
        
        return languages or ['JavaScript']  # Default to JavaScript if nothing detected
    
    def _generate_frontend_architecture_diagram(self) -> str:
        """Generate dynamic frontend architecture diagram with real data"""
        # Generate a simple text-based diagram since ProfessionalDiagramGenerator may not be available
        frontend_tech = self._repo_analysis.get('frontend_tech', [])
        if not frontend_tech:
            frontend_tech = self._extract_frontend_tech_from_files()
        
        component_names = self._extract_real_component_names()
        components_count = len(component_names) or self._repo_analysis.get('components_total', 0)
        
        # Create a simple text diagram
        tech_str = ', '.join(frontend_tech) if frontend_tech else 'React'
        
        diagram = f"""
┌─────────────────────────────────────┐
│           Frontend Stack            │
│         ({tech_str})                │
├─────────────────────────────────────┤
│  Components: {components_count:<2}                   │
│  ┌─────────────────────────────────┐│
│  │  User Interface Components      ││
│  │  • Navigation                   ││
│  │  • Forms & Inputs               ││
│  │  • Data Display                 ││
│  └─────────────────────────────────┘│
│  ┌─────────────────────────────────┐│
│  │  State Management               ││
│  │  • Application State            ││
│  │  • Component State              ││
│  └─────────────────────────────────┘│
│  ┌─────────────────────────────────┐│
│  │  API Communication              ││
│  │  • HTTP Requests                ││
│  │  • Data Fetching                ││
│  └─────────────────────────────────┘│
└─────────────────────────────────────┘"""
        
        return diagram
    
    def _extract_real_component_names(self) -> List[str]:
        """Extract real component names from repository structure"""
        component_names = []
        folder_structure = self._repo_analysis.get('folder_structure', {})
        
        for folder, info in folder_structure.items():
            if isinstance(info, dict) and 'files' in info:
                for file in info['files']:
                    # Extract component names from React/Vue/Angular files
                    if file.endswith(('.jsx', '.tsx', '.vue', '.component.ts')):
                        # Remove extension and clean name
                        name = file.replace('.jsx', '').replace('.tsx', '').replace('.vue', '').replace('.component.ts', '')
                        if name and name not in ['index', 'App', 'main']:
                            component_names.append(name)
                    elif file.endswith(('.js', '.ts')) and not file.endswith(('.test.js', '.test.ts', '.spec.js', '.spec.ts')):
                        # Check if it's likely a component file
                        name = file.replace('.js', '').replace('.ts', '')
                        if name and name[0].isupper() and name not in ['App', 'Index', 'Main']:
                            component_names.append(name)
        
        return list(set(component_names))[:10]  # Return unique names, limit to 10
    
    def _extract_frontend_tech_from_files(self) -> List[str]:
        """Extract frontend technologies from file extensions and package files"""
        tech = []
        folder_structure = self._repo_analysis.get('folder_structure', {})
        
        has_react = False
        has_vue = False
        has_angular = False
        has_typescript = False
        
        for folder, info in folder_structure.items():
            if isinstance(info, dict) and 'files' in info:
                for file in info['files']:
                    if file.endswith(('.jsx', '.tsx')):
                        has_react = True
                        if file.endswith('.tsx'):
                            has_typescript = True
                    elif file.endswith('.vue'):
                        has_vue = True
                    elif file.endswith('.component.ts'):
                        has_angular = True
                        has_typescript = True
                    elif file.endswith('.ts') and not file.endswith('.d.ts'):
                        has_typescript = True
        
        if has_react:
            tech.append('React')
        if has_vue:
            tech.append('Vue.js')
        if has_angular:
            tech.append('Angular')
        if has_typescript:
            tech.append('TypeScript')
        elif not tech:  # Default to JavaScript if no specific framework detected
            tech.append('JavaScript')
        
        return tech
    
    def _analyze_frontend_structure(self) -> List[str]:
        """Analyze frontend structure and return key insights"""
        structure_insights = []
        folder_structure = self._repo_analysis.get('folder_structure', {})
        
        # Analyze folder patterns
        frontend_folders = []
        for folder, info in folder_structure.items():
            if isinstance(info, dict) and 'files' in info:
                files = info['files']
                js_files = [f for f in files if f.endswith(('.js', '.jsx', '.ts', '.tsx', '.vue'))]
                if js_files:
                    frontend_folders.append((folder, len(js_files)))
        
        # Sort by file count
        frontend_folders.sort(key=lambda x: x[1], reverse=True)
        
        # Generate insights
        if frontend_folders:
            structure_insights.append(f"Main frontend folders: {', '.join([f[0] for f in frontend_folders[:3]])}")
        
        # Analyze file types
        total_js_files = sum(count for _, count in frontend_folders)
        if total_js_files > 0:
            structure_insights.append(f"Total frontend files: {total_js_files}")
        
        # Detect patterns
        folder_names = [folder.lower() for folder, _ in frontend_folders]
        if any('component' in name for name in folder_names):
            structure_insights.append("Component-based architecture detected")
        if any('page' in name or 'view' in name for name in folder_names):
            structure_insights.append("Page-based routing structure detected")
        if any('service' in name or 'api' in name for name in folder_names):
            structure_insights.append("Service layer for API communication detected")
        
        return structure_insights
    
    def _get_comprehensive_api_endpoints(self) -> List[Dict]:
        """Get comprehensive API endpoints from repository, PRD, and frontend analysis"""
        all_endpoints = []
        
        # 1. Repository detected endpoints
        repo_endpoints = self._repo_analysis.get('api_endpoints', [])
        for endpoint in repo_endpoints:
            if isinstance(endpoint, dict):
                all_endpoints.append(endpoint)
        
        # 2. PRD extracted endpoints (now returns structured Dict objects)
        prd_endpoints = self._extract_prd_endpoints()
        all_endpoints.extend(prd_endpoints)
        
        # 3. Analyze frontend code for API calls
        frontend_endpoints = self._analyze_frontend_for_apis()
        all_endpoints.extend(frontend_endpoints)
        
        # Remove duplicates based on path
        seen_paths = set()
        unique_endpoints = []
        for endpoint in all_endpoints:
            path = endpoint.get('path', str(endpoint))
            if path not in seen_paths:
                seen_paths.add(path)
                unique_endpoints.append(endpoint)
        
        return unique_endpoints
    
    def _analyze_frontend_for_apis(self) -> List[Dict]:
        """Analyze frontend structure to infer API endpoints from actual code patterns"""
        endpoints = []
        folder_structure = self._repo_analysis.get('folder_structure', {})
        
        # Extract API patterns from file names and folder structure
        api_indicators = set()
        
        for folder, info in folder_structure.items():
            if isinstance(info, dict) and 'files' in info:
                for file in info['files']:
                    # Extract meaningful names from frontend files
                    if file.endswith(('.js', '.jsx', '.ts', '.tsx')):
                        # Extract entity names from file names
                        file_base = file.replace('.js', '').replace('.jsx', '').replace('.ts', '').replace('.tsx', '')
                        
                        # Skip generic files
                        if file_base.lower() not in ['app', 'index', 'main', 'component', 'utils', 'config']:
                            # Clean and extract meaningful entity names
                            entity = re.sub(r'[^a-zA-Z]', '', file_base.lower())
                            if len(entity) > 2:
                                api_indicators.add(entity)
        
        # Generate endpoints from discovered entities
        for entity in api_indicators:
            endpoints.extend(self._generate_entity_endpoints(entity))
        
        return endpoints[:8]  # Limit to most relevant
    
    def _generate_backend_architecture_diagram(self, endpoints: List[Dict]) -> str:
        """Generate professional backend architecture diagram"""
        from utils.diagram_generator import ProfessionalDiagramGenerator
        
        diagram_gen = ProfessionalDiagramGenerator()
        frontend_tech = self._repo_analysis.get('frontend_tech', ['React'])
        backend_tech = self._repo_analysis.get('backend_tech', ['FastAPI'])
        database_tech = self._repo_analysis.get('database_tech', ['PostgreSQL'])
        
        # Create a mock repo analysis for the diagram generator
        mock_repo = MockRepoAnalysis({
            'frontend_tech': frontend_tech,
            'backend_tech': backend_tech,
            'database_tech': database_tech,
            'api_endpoints': endpoints
        })
        
        # Generate a simple text-based system diagram
        diagram = f"""
┌─────────────────┐    ┌─────────────────┐    ┌─────────────────┐
│    Frontend     │    │     Backend     │    │    Database     │
│  {frontend_tech[0]:<13} │◄──►│  {backend_tech[0]:<13} │◄──►│  {database_tech[0]:<13} │
│                 │    │                 │    │                 │
│ • UI Components │    │ • API Endpoints │    │ • Data Storage  │
│ • State Mgmt    │    │ • Business Logic│    │ • Transactions  │
│ • Routing       │    │ • Authentication│    │ • Indexing      │
└─────────────────┘    └─────────────────┘    └─────────────────┘
                              │
                              ▼
                    {len(endpoints)} API Endpoints
                    ┌─────────────────┐
                    │ RESTful Services│
                    │ • CRUD Ops      │
                    │ • Data Validation│
                    │ • Error Handling│
                    └─────────────────┘"""
        
        return diagram
    
    def _group_endpoints_by_service(self, endpoints: List[Dict]) -> Dict[str, List[Dict]]:
        """Group endpoints by service/domain"""
        services = {}
        
        for endpoint in endpoints:
            if isinstance(endpoint, dict):
                path = endpoint.get('path', '')
                # Extract service name from path
                if path.startswith('/api/'):
                    parts = path.split('/')
                    if len(parts) > 2:
                        service_name = parts[2].capitalize()
                    else:
                        service_name = 'General'
                else:
                    service_name = 'General'
                
                if service_name not in services:
                    services[service_name] = []
                services[service_name].append(endpoint)
        
        return services
    
    def _analyze_backend_services(self, endpoints: List[Dict]) -> List[str]:
        """Analyze backend services and return insights"""
        insights = []
        
        if not endpoints:
            insights.append("No API endpoints detected - consider implementing RESTful API")
            return insights
        
        # Analyze endpoint patterns
        methods = [ep.get('method', 'GET') for ep in endpoints if isinstance(ep, dict)]
        method_counts = {method: methods.count(method) for method in set(methods)}
        
        insights.append(f"HTTP methods distribution: {', '.join([f'{k}: {v}' for k, v in method_counts.items()])}")
        
        # Analyze service domains
        service_groups = self._group_endpoints_by_service(endpoints)
        if service_groups:
            insights.append(f"Service domains identified: {', '.join(service_groups.keys())}")
        
        # Authentication analysis
        auth_endpoints = [ep for ep in endpoints if isinstance(ep, dict) and 'auth' in ep.get('path', '').lower()]
        if auth_endpoints:
            insights.append(f"Authentication endpoints: {len(auth_endpoints)} detected")
        
        # CRUD analysis
        crud_patterns = {'GET': 'Read', 'POST': 'Create', 'PUT': 'Update', 'DELETE': 'Delete'}
        crud_ops = [crud_patterns.get(method, method) for method in methods if method in crud_patterns]
        if crud_ops:
            insights.append(f"CRUD operations supported: {', '.join(set(crud_ops))}")
        
        return insights
    
    def _infer_architecture_pattern(self, frontend_tech: List[str], backend_tech: List[str]) -> List[str]:
        """Infer architecture pattern from technology stack"""
        if frontend_tech and backend_tech:
            return ['Client-Server Architecture', 'RESTful API Design']
        elif frontend_tech:
            return ['Component-Based Architecture', 'SPA (Single Page Application)']
        elif backend_tech:
            return ['Service-Oriented Architecture', 'API-First Design']
        else:
            return ['Layered Architecture']

    def parse_prd_content(self, prd_content: str) -> Dict[str, Any]:
        """Parse PRD content to extract structured information"""
        if not prd_content:
            return {
                'product_name': 'Application',
                'tech_stack': {'languages': [], 'frontend': [], 'backend': [], 'databases': []},
                'features': [],
                'api_endpoints': [],
                'database_tables': []
            }
        
        # Clean the content first
        cleaned_content = self._clean_extracted_text(prd_content)
        
        # Extract product name using enhanced method
        product_name = self._extract_enhanced_product_name(cleaned_content)
        logger.info(f"🔍 DEBUG: Extracted product name: '{product_name}'")
        logger.info(f"🔍 DEBUG: Cleaned content preview: '{cleaned_content[:200]}...'")
        
        # Extract tech stack with enhanced patterns
        tech_stack = {'languages': [], 'frontend': [], 'backend': [], 'databases': []}
        
        # Enhanced tech stack extraction
        content_lower = cleaned_content.lower()
        
        # Define comprehensive technology mappings
        tech_keywords = {
            'languages': ['python', 'javascript', 'java', 'c#', 'csharp', 'php', 'ruby', 'go', 'rust', 'typescript', 'kotlin', 'swift', 'scala', 'clojure'],
            'frontend': ['react', 'vue', 'angular', 'next.js', 'nextjs', 'nuxt', 'svelte', 'html', 'css', 'bootstrap', 'tailwind', 'material-ui', 'chakra'],
            'backend': ['fastapi', 'flask', 'django', 'express', 'spring', 'laravel', 'rails', 'gin', 'actix', 'node.js', 'nodejs', 'asp.net', 'koa'],
            'databases': ['postgresql', 'postgres', 'mysql', 'mongodb', 'redis', 'sqlite', 'oracle', 'cassandra', 'dynamodb', 'elasticsearch', 'mariadb']
        }
        
        # Extract technologies from content
        for category, keywords in tech_keywords.items():
            for keyword in keywords:
                if keyword in content_lower:
                    # Normalize names
                    display_name = {
                        'nextjs': 'Next.js', 'nodejs': 'Node.js', 'csharp': 'C#',
                        'postgres': 'PostgreSQL', 'mongodb': 'MongoDB'
                    }.get(keyword, keyword.title())
                    
                    if display_name not in tech_stack[category]:
                        tech_stack[category].append(display_name)
        
        # Extract features with enhanced patterns
        features = []
        
        # Look for feature sections
        feature_sections = re.findall(r'(?:core features?|main features?|key features?|features?|functionality)[:\s]*\n([\s\S]*?)(?=\n\d+\.|\n[A-Z][^\n]*:|$)', cleaned_content, re.IGNORECASE)
        
        for section in feature_sections:
            # Extract bullet points and numbered items
            feature_items = re.findall(r'(?:^|\n)\s*(?:[\-•*]|\d+\.?)\s*([^\n]+)', section, re.MULTILINE)
            for item in feature_items:
                clean_feature = item.strip()
                if len(clean_feature) > 5 and len(clean_feature) < 100:
                    features.append(clean_feature)
        
        # Also look for "As a user" stories
        user_stories = re.findall(r'as a (?:user|admin|customer)[^\n]*?(?:want to|can)\s+([^\n\.]+)', cleaned_content, re.IGNORECASE)
        for story in user_stories:
            clean_story = story.strip()
            if len(clean_story) > 5 and len(clean_story) < 100:
                features.append(clean_story)
        
        # Extract API endpoints with comprehensive patterns
        api_endpoints = []
        
        # Look for API endpoint sections
        api_sections = re.findall(r'(?:api endpoints?|endpoints?)[:\s]*\n([\s\S]*?)(?=\n\d+\.|\n[A-Z][^\n]*:|$)', cleaned_content, re.IGNORECASE)
        
        for section in api_sections:
            # Extract endpoint patterns
            endpoint_patterns = [
                r'/[a-zA-Z][a-zA-Z0-9/_-]*',  # Path patterns
                r'(?:^|\n)\s*(?:[\-•*]|\d+\.?)\s*([^\n]*(?:api|endpoint)[^\n]*)',  # Bullet points with API/endpoint
                r'(GET|POST|PUT|DELETE|PATCH)\s+(/[^\s\n]+)',  # HTTP methods with paths
            ]
            
            for pattern in endpoint_patterns:
                matches = re.findall(pattern, section, re.IGNORECASE | re.MULTILINE)
                for match in matches:
                    if isinstance(match, tuple):
                        api_endpoints.append(f"{match[0]} {match[1]}")
                    elif isinstance(match, str) and len(match.strip()) > 3:
                        api_endpoints.append(match.strip())
        
        # Extract database tables
        database_tables = []
        
        # Look for database schema sections
        db_sections = re.findall(r'(?:database schema|schema|tables?)[:\s]*\n([\s\S]*?)(?=\n\d+\.|\n[A-Z][^\n]*:|$)', cleaned_content, re.IGNORECASE)
        
        for section in db_sections:
            # Extract table names from various formats
            table_patterns = [
                r'(?:^|\n)\s*(?:[\-•*]|\d+\.?)\s*([A-Za-z][A-Za-z0-9_]*)',  # Bullet points with table names
                r'\b([A-Za-z][A-Za-z0-9_]*)\s+(?:table|entity)',  # Table/entity mentions
            ]
            
            for pattern in table_patterns:
                matches = re.findall(pattern, section, re.MULTILINE)
                for match in matches:
                    clean_table = match.strip().lower()
                    if len(clean_table) > 2 and clean_table not in database_tables:
                        database_tables.append(clean_table)
        
        # If no explicit tables found, infer from content
        if not database_tables:
            common_entities = ['users', 's', 'rooms', 'bookings', 'products', 'orders', 'customers']
            for entity in common_entities:
                if entity.rstrip('s') in content_lower:  # Check singular form
                    database_tables.append(entity)
        
        result = {
            'product_name': product_name,
            'tech_stack': tech_stack,
            'features': features[:15],  # Limit features
            'api_endpoints': api_endpoints[:20],  # Limit endpoints
            'database_tables': database_tables[:10],  # Limit tables
            'content': cleaned_content
        }
        
        logger.info(f"🔍 DEBUG: Final PRD parse result - Product: '{result['product_name']}'")
        return result
    
    def _preprocess_prd_content(self, content: str) -> str:
        """Preprocess PRD content to improve parsing accuracy"""
        # Remove excessive whitespace and normalize line endings
        content = re.sub(r'\r\n', '\n', content)
        content = re.sub(r'\r', '\n', content)
        content = re.sub(r'\n\s*\n', '\n\n', content)
        
        # Remove PDF artifacts and metadata
        content = re.sub(r'%PDF-[\d\.]+.*?\n', '', content, flags=re.DOTALL)
        content = re.sub(r'/[A-Z][a-zA-Z]*\s+\d+\s+\d+\s+R', '', content)
        content = re.sub(r'<<.*?>>', '', content, flags=re.DOTALL)
        content = re.sub(r'stream\s*.*?\s*endstream', '', content, flags=re.DOTALL)
        
        # Clean up common document artifacts
        content = re.sub(r'Page \d+ of \d+', '', content)
        content = re.sub(r'\f', '\n', content)  # Form feed to newline
        content = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]', '', content)  # Remove control characters
        
        # Normalize bullet points and numbering
        content = re.sub(r'•', '•', content)  # Normalize bullet points
        content = re.sub(r'[\u2010-\u2015]', '-', content)  # Normalize dashes
        
        return content.strip()
    
    def _extract_enhanced_features(self, content: str) -> List[str]:
        """Enhanced feature extraction with multiple strategies"""
        features = []
        
        # Strategy 1: Look for explicit feature sections
        feature_patterns = [
            r'(?:features?|functionality|capabilities)[:\s]*\n([\s\S]*?)(?:\n\n|\n[A-Z]|$)',
            r'(?:key|core|main)\s+features?[:\s]*\n([\s\S]*?)(?:\n\n|\n[A-Z]|$)',
            r'(?:requirements?|specifications?)[:\s]*\n([\s\S]*?)(?:\n\n|\n[A-Z]|$)'
        ]
        
        for pattern in feature_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                # Extract bullet points or numbered items
                items = re.findall(r'(?:•|\*|-|\d+\.)\s*([^\n]+)', match)
                features.extend([item.strip() for item in items if len(item.strip()) > 5])
        
        # Strategy 2: Look for action verbs (user stories)
        action_patterns = [
            r'(?:user can|users can|ability to|should be able to|will be able to)\s+([^\n\.]+)',
            r'(?:the system|application|app)\s+(?:will|should|must|can)\s+([^\n\.]+)'
        ]
        
        for pattern in action_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            features.extend([match.strip() for match in matches if len(match.strip()) > 10])
        
        return list(set(features))[:15]  # Limit and deduplicate
    
    def _extract_enhanced_api_methods(self, content: str) -> List[str]:
        """Enhanced API method extraction with comprehensive patterns"""
        api_methods = []
        
        logger.info("🔍 Starting enhanced API method extraction...")
        
        # Strategy 1: HTTP methods with paths
        http_patterns = [
            r'(GET|POST|PUT|DELETE|PATCH)\s+(/[^\s\n]+)',
            r'(GET|POST|PUT|DELETE|PATCH)[:\s]+([^\n]+)',
            r'\b(GET|POST|PUT|DELETE|PATCH)\b[^\n]*(/api/[^\s\n]+)',
        ]
        
        for i, pattern in enumerate(http_patterns):
            matches = re.findall(pattern, content, re.IGNORECASE)
            logger.info(f"🔍 HTTP pattern {i+1}: Found {len(matches)} matches")
            for match in matches:
                api_method = f"{match[0]} {match[1]}"
                api_methods.append(api_method)
                logger.info(f"✅ Found HTTP method: {api_method}")
        
        # Strategy 2: API endpoint descriptions and paths
        endpoint_patterns = [
            r'(?:endpoint|api)[:\s]*([^\n]+)',
            r'/api/[a-zA-Z0-9/_-]+',
            r'(?:route|path)[:\s]*(/[^\s\n]+)',
            r'\b(\w+)\s+endpoint',
            r'\b(\w+)\s+API',
            r'•\s*([^\n]*(?:endpoint|API|api)[^\n]*)',
            r'-\s*([^\n]*(?:endpoint|API|api)[^\n]*)',
            r'\d+\.\s*([^\n]*(?:endpoint|API|api)[^\n]*)',
        ]
        
        for i, pattern in enumerate(endpoint_patterns):
            matches = re.findall(pattern, content, re.IGNORECASE)
            logger.info(f"🔍 Endpoint pattern {i+1}: Found {len(matches)} matches")
            for match in matches:
                if isinstance(match, str) and match.strip():
                    api_methods.append(match.strip())
                    logger.info(f"✅ Found endpoint: {match.strip()}")
        
        # Strategy 3: Common API functionality keywords
        api_keywords = [
            'login', 'register', 'authenticate', 'logout',
            'user', 'profile', 'account',
            'product', 'item', 'catalog',
            'order', 'purchase', 'payment',
            'cart', 'basket', 'checkout',
            'search', 'filter', 'query',
            'upload', 'download', 'file',
            'data', 'information', 'record'
        ]
        
        content_lower = content.lower()
        for keyword in api_keywords:
            if keyword in content_lower:
                api_methods.append(f"API for {keyword} functionality")
                logger.info(f"✅ Inferred API from keyword '{keyword}': API for {keyword} functionality")
        
        # Strategy 4: Look for numbered or bulleted lists that might contain endpoints
        list_patterns = [
            r'\d+\.\s*([^\n]+)',  # Numbered lists
            r'•\s*([^\n]+)',    # Bullet points
            r'-\s*([^\n]+)',      # Dash points
            r'\*\s*([^\n]+)',     # Asterisk points
        ]
        
        for pattern in list_patterns:
            matches = re.findall(pattern, content, re.MULTILINE)
            for match in matches:
                if any(keyword in match.lower() for keyword in ['api', 'endpoint', 'service', 'function', 'method']):
                    api_methods.append(match.strip())
                    logger.info(f"✅ Found from list: {match.strip()}")
        
        unique_methods = list(set(api_methods))[:25]  # Increased limit
        logger.info(f"🎯 Total unique API methods extracted: {len(unique_methods)}")
        
        return unique_methods
    
    def _extract_enhanced_database_tables(self, content: str) -> List[str]:
        """Enhanced database table extraction"""
        tables = []
        
        # Strategy 1: Explicit table mentions
        table_patterns = [
            r'(?:table|entity|model)[:\s]*([a-zA-Z_][a-zA-Z0-9_]*)',
            r'(?:database|db)\s+(?:table|schema)[:\s]*([a-zA-Z_][a-zA-Z0-9_]*)',
            r'CREATE\s+TABLE\s+([a-zA-Z_][a-zA-Z0-9_]*)',
        ]
        
        for pattern in table_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            tables.extend([match.strip() for match in matches if len(match.strip()) > 2])
        
        # Strategy 2: Infer from features (common entities)
        content_lower = content.lower()
        common_entities = ['user', 'product', 'order', 'customer', 'item', 'category', 'payment', 'account']
        
        for entity in common_entities:
            if entity in content_lower:
                tables.append(f"{entity}s")  # Pluralize
        
        return list(set(tables))[:10]  # Limit and deduplicate
    
    def _extract_enhanced_goals(self, content: str) -> List[str]:
        """Enhanced goals extraction"""
        goals = []
        
        # Strategy 1: Explicit goal sections
        goal_patterns = [
            r'(?:goal|objective|purpose|aim)[s]?[:\s]*\n([\s\S]*?)(?:\n\n|\n[A-Z]|$)',
            r'(?:mission|vision)[:\s]*([^\n\.]+)'
        ]
        
        for pattern in goal_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                # Extract bullet points or sentences
                items = re.findall(r'(?:•|\*|-|\d+\.)\s*([^\n]+)', match)
                if items:
                    goals.extend([item.strip() for item in items if len(item.strip()) > 10])
                else:
                    # Single sentence goal
                    if len(match.strip()) > 10:
                        goals.append(match.strip())
        
        return list(set(goals))[:10]  # Limit and deduplicate
    
    def _apply_intelligent_defaults(self, content: str, prd_analysis: Dict) -> Dict:
        """Apply intelligent defaults when PRD parsing yields no results"""
        
        content_lower = content.lower()
        
        # Default features based on content keywords
        if 'user' in content_lower:
            prd_analysis['features'].append('User management and authentication')
        if 'data' in content_lower:
            prd_analysis['features'].append('Data storage and retrieval')
        if 'api' in content_lower:
            prd_analysis['features'].append('RESTful API endpoints')
        if 'web' in content_lower or 'app' in content_lower:
            prd_analysis['features'].append('Web application interface')
        
        # Default API methods
        if not prd_analysis['api_methods']:
            prd_analysis['api_methods'] = [
                'GET /api/health',
                'POST /api/auth/login',
                'GET /api/data',
                'POST /api/data'
            ]
        
        # Default database tables
        if not prd_analysis['database_tables']:
            prd_analysis['database_tables'] = ['users', 'data', 'sessions']
        
        # Default goals
        if not prd_analysis['goals']:
            prd_analysis['goals'] = [
                'Provide reliable and scalable application',
                'Ensure secure user data management',
                'Deliver intuitive user experience'
            ]
        
        return prd_analysis
    
    def _infer_from_content(self, content: str, base_analysis: Dict) -> Dict:
        """Extract actual content without hardcoded assumptions"""
        # Only extract what's actually in the content
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        # Extract any bullet points or numbered items as features
        extracted_features = []
        extracted_apis = []
        extracted_tables = []
        
        for line in lines:
            line_clean = line.strip()
            # Look for actual bullet points or numbered items
            if re.match(r'^[\u2022\*\-\d+\.]\s+', line_clean):
                clean_item = re.sub(r'^[\u2022\*\-\d+\.\s]+', '', line_clean)
                if len(clean_item) > 3:
                    if any(word in clean_item.lower() for word in ['api', 'endpoint', 'post', 'get', 'put', 'delete']):
                        extracted_apis.append(clean_item)
                    elif any(word in clean_item.lower() for word in ['table', 'model', 'entity', 'schema']):
                        extracted_tables.append(clean_item)
                    else:
                        extracted_features.append(clean_item)
        
        # Only update if we found actual content
        if extracted_features:
            base_analysis['features'] = extracted_features[:10]
        if extracted_apis:
            base_analysis['api_methods'] = extracted_apis[:8]
        if extracted_tables:
            base_analysis['database_tables'] = extracted_tables[:5]
        
        return base_analysis
    
    def _get_default_prd_analysis(self) -> Dict:
        """Return intelligent default analysis when PRD parsing fails"""
        return {
            'product_name': 'Web Application',
            'goals': [
                'Deliver scalable and maintainable application',
                'Provide secure user authentication',
                'Ensure optimal user experience'
            ],
            'features': [
                'User authentication and authorization',
                'Data management and storage',
                'RESTful API endpoints',
                'Responsive web interface'
            ],
            'tech_stack': ['Python', 'FastAPI', 'PostgreSQL'],
            'database_tables': ['users', 'sessions', 'data'],
            'api_methods': [
                'GET /api/health',
                'POST /api/auth/login',
                'GET /api/users/profile',
                'POST /api/data',
                'GET /api/data'
            ],
            'content': 'Default PRD content for web application'
        }

    def _get_dynamic_report_title(self) -> str:
        """Generate dynamic report title from project name"""
        project_name = self._extract_project_title()
        return f"{project_name} - Architecture Report"
    
    def _extract_project_title(self) -> str:
        """Extract project title from GitHub URL or PRD document"""
        # Priority 1: Extract from PRD document
        prd_title = self._extract_title_from_prd()
        if prd_title and prd_title not in ['Application', 'Web Application', 'Python Desktop Application']:
            return prd_title
        
        # Priority 2: Extract from GitHub URL
        github_title = self._extract_title_from_github()
        if github_title:
            return github_title
        
        # Fallback
        return "System Architecture Report"
    
    def _extract_title_from_prd(self) -> str:
        """Extract actual title from PRD content without hardcoded keywords"""
        prd_content = self._prd_analysis.get('content', '')
        if not prd_content:
            return ''
        
        lines = [line.strip() for line in prd_content.split('\n')[:15] if line.strip()]
        
        for line in lines:
            clean_line = re.sub(r'^[#*\-•\s\d\.]+', '', line).strip()
            
            # Skip common document headers
            if re.match(r'^(table of contents?|prepared by|document|version|date|author)', clean_line, re.IGNORECASE):
                continue
            
            # Look for meaningful titles (not too short, not too long, contains letters)
            if (len(clean_line) > 8 and len(clean_line) < 100 and 
                re.search(r'[a-zA-Z]', clean_line) and
                not clean_line.lower().startswith(('by:', 'author:', 'date:', 'version:'))):
                
                # Clean up common PRD artifacts
                title = re.sub(r'\b(prd|product requirements? document|requirements?)\b', '', clean_line, flags=re.IGNORECASE).strip()
                title = re.sub(r'\s+', ' ', title).strip()
                
                if title and len(title) > 5:
                    return title
        
        return ''
    
    def _extract_title_from_github(self) -> str:
        """Extract title from GitHub repository URL"""
        github_url = getattr(self, '_github_url', '')
        if not github_url:
            return ''
        
        # Extract repo name from GitHub URL
        match = re.search(r'github\.com/[^/]+/([^/\.]+)', github_url)
        if match:
            repo_name = match.group(1)
            # Convert repo name to readable title
            title = re.sub(r'[_-]', ' ', repo_name)
            
            # Capitalize words properly
            words = []
            for word in title.split():
                if len(word) > 1:  # Keep words longer than 1 character
                    words.append(word.capitalize())
            
            if words:
                return ' '.join(words)
        
        return ''

    def generate_architecture_pdf(self, architecture, github_url="", prd_included=False, repo_analysis=None, prd_content=None, prd_file_path=None) -> str:
        """Generate PDF with 100% dynamic content"""
        
        # Store GitHub URL for title extraction
        self._github_url = github_url
        
        # Handle repo_analysis object conversion
        if repo_analysis and hasattr(repo_analysis, '__dict__'):
            self._repo_analysis = self.analyze_repo_from_object(repo_analysis)
        else:
            self._repo_analysis = {}
        
        # Handle PRD content - support file path or direct content
        if prd_content:
            self._prd_analysis = self.parse_prd_content(prd_content)
        elif prd_file_path and os.path.exists(prd_file_path):
            try:
                prd_text = self.extract_text_from_file(prd_file_path)
                self._prd_analysis = self.parse_prd_content(prd_text)
                logger.info(f"Successfully extracted PRD from file: {prd_file_path}")
            except Exception as e:
                logger.error(f"Failed to extract PRD from file {prd_file_path}: {str(e)}")
                self._prd_analysis = self._get_default_prd_analysis()
        else:
            self._prd_analysis = self._get_default_prd_analysis()
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        product_name = self._prd_analysis.get('product_name', 'Application').replace(' ', '_').replace('/', '_')
        # Use safe filename
        safe_filename = re.sub(r'[^\w\-_\.]', '_', f"architecture_report_{product_name}_{timestamp}")
        filename = f"{safe_filename}.pdf"
        filepath = os.path.join(self.output_dir, filename)
        
        # Ensure output directory exists and is writable
        os.makedirs(self.output_dir, exist_ok=True)
        
        doc = SimpleDocTemplate(filepath, pagesize=A4, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
        story = []
        
        # Cover Page - Dynamic
        story.append(Paragraph(self._get_dynamic_report_title(), self.styles['CustomTitle']))
        story.append(Spacer(1, 0.3*inch))
        story.append(Paragraph(self._get_product_name(), self.styles['CustomHeading1']))
        story.append(Spacer(1, 0.2*inch))
        story.append(Paragraph(self._get_description(), self.styles['CustomBody']))
        story.append(Spacer(1, 0.3*inch))
        story.append(self._create_cover_table(github_url, prd_included))
        story.append(Spacer(1, 0.2*inch))
        story.append(self._create_stats_table())
        story.append(PageBreak())
        
        # TOC
        story.append(Paragraph("Table of Contents", self.styles['CustomHeading1']))
        
        # Add each TOC item as separate paragraph for proper vertical formatting
        toc_items = self._get_toc().split('\n')
        for item in toc_items:
            story.append(Paragraph(item, self.styles['CustomBody']))
        
        story.append(PageBreak())
        
        # All sections with dynamic content
        story.extend(self._create_executive_summary())
        story.append(PageBreak())
        story.extend(self._create_goals_scope())
        story.append(PageBreak())
        story.extend(self._create_context_diagram())
        story.append(PageBreak())
        story.extend(self._create_frontend_section())
        story.append(PageBreak())
        story.extend(self._create_backend_section())
        story.append(PageBreak())
        story.extend(self._create_api_section())
        story.append(PageBreak())
        story.extend(self._create_sequence_diagram_section())
        story.append(PageBreak())
        story.extend(self._create_interactions_section())
        story.append(PageBreak())
        story.extend(self._create_unified_diagram())
        story.append(PageBreak())
        story.extend(self._create_deployment_section())
        story.append(PageBreak())
        story.extend(self._create_security_section())
        story.append(PageBreak())
        story.extend(self._create_tech_stack_section())
        story.append(PageBreak())
        if prd_included:
            story.extend(self._create_business_alignment())
            story.append(PageBreak())
        story.extend(self._create_recommendations_section())
        story.append(PageBreak())
        
        # Add system architecture diagram
        story.extend(self._create_architecture_diagram_section())
        story.append(PageBreak())
        
        # Add layered data flow diagram
        story.extend(self._create_layered_dataflow_section())
        
        doc.build(story)
        return filepath

    # Dynamic content methods
    def _get_product_name(self) -> str:
        product_name = self._prd_analysis.get('product_name', 'Application')
        logger.info(f"🔍 DEBUG: Raw product_name from PRD analysis: '{product_name}'")
        logger.info(f"🔍 DEBUG: PRD analysis keys: {list(self._prd_analysis.keys()) if self._prd_analysis else 'None'}")
        
        # Clean up corrupted product names
        if product_name.startswith('%PDF') or 'PDF-1.' in product_name:
            logger.info(f"🔍 DEBUG: Detected corrupted PDF product name, using fallback")
            return 'Web Application'
        
        # Always use the extracted product name if it's not the default
        if product_name and product_name != 'Application':
            logger.info(f"🔍 DEBUG: Using extracted product name: '{product_name}'")
            return product_name
        
        logger.info(f"🔍 DEBUG: Using default 'Application' name")
        return 'Application'
    
    def _extract_enhanced_product_name(self, prd_content: str) -> str:
        """Enhanced product name extraction from PRD content"""
        if not prd_content:
            logger.info(f"🔍 DEBUG: No PRD content provided")
            return "Application"
        
        lines = prd_content.strip().split('\n')
        logger.info(f"🔍 DEBUG: Processing {len(lines)} lines for product name")
        
        # Look for title patterns in first 10 lines
        for i, line in enumerate(lines[:10]):
            cleaned_line = line.strip()
            logger.info(f"🔍 DEBUG: Line {i+1}: '{cleaned_line[:100]}...'")
            
            # Skip empty lines
            if not cleaned_line:
                continue
            
            # Remove common prefixes and clean the line
            cleaned_line = re.sub(r'^[#*\-•\s]+', '', cleaned_line)  # Remove markdown, bullets
            cleaned_line = re.sub(r'[""''„"«»]', '"', cleaned_line)  # Normalize quotes
            cleaned_line = re.sub(r'[–—−]', '-', cleaned_line)  # Normalize dashes
            cleaned_line = re.sub(r'[\x00-\x1f\x7f-\x9f]', ' ', cleaned_line)  # Remove control chars
            cleaned_line = re.sub(r'\s+', ' ', cleaned_line).strip()  # Normalize whitespace
            
            logger.info(f"🔍 DEBUG: Cleaned line {i+1}: '{cleaned_line}'")
            
            # Skip metadata lines
            if re.match(r'^(Page|Document|Version|Prepared|Table|Product Requirements|PRD|Contents?)\b', cleaned_line, re.IGNORECASE):
                logger.info(f"🔍 DEBUG: Skipping metadata line: '{cleaned_line}'")
                continue
            
            # Look for application/system/platform keywords that indicate a title
            title_keywords = ['application', 'system', 'platform', 'app', 'service', 'portal', 'dashboard', 'management', 'booking', 'ecommerce', 'marketplace', 'api', 'website', 'tool', 'solution']
            
            # Check if this looks like a title
            has_keyword = any(keyword in cleaned_line.lower() for keyword in title_keywords)
            matches_pattern = re.match(r'^[A-Z][a-zA-Z\s]+[A-Za-z]$', cleaned_line)
            
            logger.info(f"🔍 DEBUG: Line '{cleaned_line}' - Length: {len(cleaned_line)}, Has keyword: {has_keyword}, Matches pattern: {bool(matches_pattern)}")
            
            if (len(cleaned_line) > 3 and 
                len(cleaned_line) < 100 and  # Reasonable title length
                not re.match(r'^\d+\.?\s*$', cleaned_line) and  # Not just numbers
                not cleaned_line.lower().startswith(('by:', 'author:', 'date:', 'overview')) and
                (has_keyword or matches_pattern)):
                logger.info(f"🔍 DEBUG: FOUND PRODUCT NAME: '{cleaned_line}'")
                return cleaned_line
        
        logger.info(f"🔍 DEBUG: No product name found, using default 'Application'")
        return "Application"

    def _get_description(self) -> str:
        product = self._get_product_name()
        has_frontend = bool(self._repo_analysis.get('frontend_tech'))
        has_backend = bool(self._repo_analysis.get('backend_tech'))
        
        if has_frontend and has_backend:
            app_type = "full-stack web application"
        elif has_frontend:
            app_type = "frontend UI application"
        elif has_backend:
            app_type = "backend service"
        else:
            app_type = "web application"
        
        tech_summary = self._get_tech_summary()
        backend_status = 'no backend detected in repo' if not has_backend else 'backend implemented'
        return f"The {product} is a {app_type} {tech_summary}. The repository provides UI/code components. PRD specifies full features and architecture. ({backend_status})"

    def _get_tech_summary(self) -> str:
        frontend = self._repo_analysis.get('frontend_tech', [])
        backend = self._repo_analysis.get('backend_tech', [])
        
        if frontend and backend:
            return f"built with {', '.join(frontend[:2])} frontend and {', '.join(backend[:2])} backend"
        elif frontend:
            return f"built with {', '.join(frontend[:2])} for the frontend"
        elif backend:
            return f"built with {', '.join(backend[:2])} for the backend"
        else:
            # Infer from languages if available
            languages = list(self._repo_analysis.get('languages', {}).keys())
            if languages:
                return f"using {', '.join(languages[:2])} technologies"
            return "with modern web technologies"

    def _get_arch_pattern(self) -> str:
        patterns = self._repo_analysis.get('patterns', [])
        if patterns:
            return ', '.join(patterns)
        elif self._repo_analysis.get('frontend_tech') and self._repo_analysis.get('backend_tech'):
            return "Client-Server Architecture"
        else:
            return "Layered Architecture"

    def _get_complexity_score(self) -> int:
        base = len(self._repo_analysis.get('languages', {})) * 2
        base += len(self._repo_analysis.get('api_endpoints', [])) // 5
        base += self._repo_analysis.get('components_total', 0) // 10
        return min(10, max(3, base))

    def _get_scalability_level(self) -> str:
        if any('redis' in tech.lower() for tech in self._repo_analysis.get('database_tech', [])):
            return "High Scalability"
        elif self._repo_analysis.get('api_endpoints'):
            return "Moderate Scalability"
        else:
            return "Basic Scalability"

    def _create_cover_table(self, github_url: str, prd_included: bool) -> Table:
        file_count = self._repo_analysis.get('file_count', 0)
        
        # Determine analysis type based on actual detected technologies
        frontend_tech = self._repo_analysis.get('frontend_tech', [])
        backend_tech = self._repo_analysis.get('backend_tech', [])
        
        if frontend_tech and backend_tech:
            analysis_type = 'full-stack analysis'
        elif frontend_tech and not backend_tech:
            analysis_type = 'frontend only'
        elif backend_tech and not frontend_tech:
            analysis_type = 'backend only'
        elif file_count > 0:
            analysis_type = 'code analysis'
        else:
            analysis_type = 'no files detected'
        
        # Show actual file count with better fallback
        if file_count > 0:
            scope_text = f"{file_count} files analyzed ({analysis_type})"
        else:
            # Check if we have any other indicators of repository content
            has_languages = bool(self._repo_analysis.get('languages'))
            has_components = self._repo_analysis.get('components_total', 0) > 0
            has_tech = bool(self._repo_analysis.get('frontend_tech')) or bool(self._repo_analysis.get('backend_tech'))
            
            if has_languages or has_components or has_tech:
                scope_text = f"Repository analyzed ({analysis_type}) - file count unavailable"
            else:
                scope_text = "No files detected in repository"
        
        data = [
            ['Generated from:', 'GitHub Repository Analysis'],
            ['Repository URL:', github_url or 'Local Analysis'],
            ['PRD Analysis:', 'Included' if prd_included else 'Not Provided'],
            ['Generated on:', datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
            ['Analysis Scope:', scope_text]
        ]
        
        table = Table(data, colWidths=[2*inch, 4*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), self.colors['light_gray']),
            ('TEXTCOLOR', (0, 0), (-1, -1), self.colors['text']),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('GRID', (0, 0), (-1, -1), 1, self.colors['medium_gray']),
        ]))
        
        return table

    def _create_stats_table(self) -> Table:
        data = [
            ['Architecture Pattern:', self._get_arch_pattern()],
            ['Application Type:', 'Web Application' if self._repo_analysis.get('frontend_tech') else 'Backend Service'],
            ['Complexity Score:', f"{self._get_complexity_score()}/10"],
            ['Scalability Level:', self._get_scalability_level()],
            ['Technology Maturity:', 'Modern' if self._repo_analysis.get('build_tools') else 'Standard']
        ]
        
        table = Table(data, colWidths=[2*inch, 4*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), self.colors['primary']),
            ('TEXTCOLOR', (0, 0), (0, -1), white),
            ('TEXTCOLOR', (1, 0), (1, -1), self.colors['text']),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('GRID', (0, 0), (-1, -1), 1, self.colors['medium_gray']),
        ]))
        
        return table

    def _get_toc(self) -> str:
        toc_items = [
            "1. Executive Summary",
            "2. Architecture Goals & Scope", 
            "3. Architecture Context Diagram",
            "4. Frontend Architecture",
            "5. Backend Architecture",
            "6. API Endpoints Documentation",
            "7. System Interaction Flow (Sequence Diagram)",
            "8. Component Interactions",
            "9. Unified Architecture Diagram",
            "10. Deployment Architecture",
            "11. Security Model",
            "12. Technology Stack Summary",
            "13. Business Alignment",
            "14. Recommendations & Next Steps",
            "15. System Architecture Diagram"
        ]
        return "\n".join(toc_items)

    def _create_executive_summary(self) -> List:
        story = []
        story.append(Paragraph("1. Executive Summary", self.styles['CustomHeading1']))
        
        # Debug what we actually have
        logger.info(f"🔍 DEBUG - repo_analysis keys: {list(self._repo_analysis.keys()) if self._repo_analysis else 'None'}")
        logger.info(f"🔍 DEBUG - prd_analysis keys: {list(self._prd_analysis.keys()) if self._prd_analysis else 'None'}")
        logger.info(f"🔍 DEBUG - repo_analysis content: {self._repo_analysis}")
        logger.info(f"🔍 DEBUG - prd_analysis content: {self._prd_analysis}")
        
        # Get actual detected values with intelligent inference
        prd_endpoints = self._extract_prd_endpoints()
        repo_endpoints = self._repo_analysis.get('api_endpoints', [])
        api_count = len(prd_endpoints) + len(repo_endpoints)
        
        # If PRD gives zero APIs, use frontend/LLM analysis as fallback
        if api_count == 0:
            api_count = self._infer_api_count_from_frontend_and_llm()
        
        frontend_count = self._detect_repo_components()
        backend_tech = self._repo_analysis.get('backend_tech', [])
        backend_count = len(backend_tech)
        
        # Infer backend technologies dynamically if none detected
        if backend_count == 0 and (api_count > 0 or frontend_count > 0):
            inferred_backend = self._infer_backend_technologies()
            backend_count = len(inferred_backend)
        
        languages = list(self._repo_analysis.get('languages', {}).keys())
        language_count = len(languages)
        
        # Infer languages dynamically if none detected
        if language_count == 0 and (frontend_count > 0 or backend_count > 0):
            inferred_languages = self._infer_programming_languages()
            language_count = len(inferred_languages)
        
        logger.info(f"🔍 RAW COUNTS - APIs: {api_count}, Frontend: {frontend_count}, Backend: {backend_count}, Languages: {language_count}")
        logger.info(f"🔍 PRD endpoints: {prd_endpoints}")
        logger.info(f"🔍 Repo endpoints: {repo_endpoints}")
        logger.info(f"🔍 Backend tech: {backend_tech}")
        logger.info(f"🔍 Languages: {languages}")
        
        # Executive summary text
        summary_text = f"This document presents a comprehensive analysis of the {self._get_product_name()} architecture, generated through automated analysis of the GitHub repository and associated documentation."
        if backend_count == 0:
            summary_text += " No backend code detected – recommendations based on PRD requirements."
        else:
            summary_text += " Backend components detected and analyzed."
        story.append(Paragraph(summary_text, self.styles['CustomBody']))
        story.append(Spacer(1, 0.2*inch))
        
        # Key findings with actual detected values
        story.append(Paragraph("Key Findings:", self.styles['CustomHeading2']))
        
        findings = [
            f"• {api_count} API endpoints (repo + PRD) identified and documented",
            f"• {frontend_count} frontend components analyzed",
            f"• {backend_count} backend technologies detected",
            f"• {language_count} programming languages in use",
            f"• {len(self._get_recommendations())} architectural recommendations provided"
        ]
        
        for finding in findings:
            story.append(Paragraph(finding, self.styles['CustomBullet']))
        
        return story
    
    def _infer_api_count_from_frontend_and_llm(self) -> int:
        """Infer API count from frontend analysis and LLM reasoning"""
        # Check if we have frontend components that would need APIs
        frontend_count = self._detect_repo_components()
        
        # Basic inference: if we have frontend components, assume some APIs
        if frontend_count > 0:
            return max(3, frontend_count // 3)  # Rough estimate
        
        # Check PRD content for API indicators
        prd_content = self._prd_analysis.get('content', '').lower()
        api_indicators = ['api', 'endpoint', 'rest', 'service', 'backend']
        
        if any(indicator in prd_content for indicator in api_indicators):
            return 5  # Default minimum for API-mentioned projects
        
        return 0
    
    def _get_enhanced_metrics(self) -> Dict[str, int]:
        """Get enhanced metrics with intelligent detection and defaults"""
        
        logger.info("📊 Starting enhanced metrics calculation...")
        
        # Force extract PRD endpoints first
        prd_endpoints = self._extract_prd_endpoints()
        logger.info(f"📊 PRD endpoints extracted: {len(prd_endpoints)}")
        
        # Force extract repository data
        repo_languages = list(self._repo_analysis.get('languages', {}).keys())
        repo_components = self._detect_repo_components()
        repo_backend_tech = self._detect_repo_backend_tech()
        repo_endpoints = self._repo_analysis.get('api_endpoints', [])
        
        logger.info(f"📊 Repo data - Languages: {repo_languages}, Components: {repo_components}, Backend: {repo_backend_tech}")
        
        # Calculate metrics with minimum guarantees
        api_count = len(repo_endpoints) + len(prd_endpoints)
        if api_count == 0:
            # Force minimum based on PRD content
            if self._prd_analysis.get('content') and len(self._prd_analysis.get('content', '')) > 100:
                api_count = max(5, len(self._prd_analysis.get('api_methods', [])))  # Minimum 5 if PRD exists
                logger.info(f"📊 Applied PRD-based API minimum: {api_count}")
        
        frontend_count = max(repo_components, self._repo_analysis.get('components_total', 0))
        if frontend_count == 0:
            # Check if repo has frontend files
            folder_structure = self._repo_analysis.get('folder_structure', {})
            frontend_files = 0
            for folder, info in folder_structure.items():
                if isinstance(info, dict) and 'files' in info:
                    for file in info['files']:
                        if file.endswith(('.js', '.jsx', '.ts', '.tsx', '.vue', '.html')):
                            frontend_files += 1
            frontend_count = max(frontend_files, 3 if frontend_files > 0 else 0)
            logger.info(f"📊 Calculated frontend count from files: {frontend_count}")
        
        backend_count = len(repo_backend_tech)
        if backend_count == 0 and api_count > 0:
            backend_count = 3  # If we have APIs, assume backend exists
            logger.info(f"📊 Applied backend default due to APIs: {backend_count}")
        
        language_count = len(repo_languages)
        if language_count == 0:
            # Infer from file extensions in repo
            detected_langs = self._detect_repo_languages()
            language_count = max(len(detected_langs), 2 if frontend_count > 0 or backend_count > 0 else 0)
            logger.info(f"📊 Detected languages: {detected_langs}, count: {language_count}")
        
        final_metrics = {
            'frontend_count': frontend_count,
            'backend_count': backend_count,
            'api_count': api_count,
            'language_count': language_count
        }
        
        logger.info(f"🎯 Final enhanced metrics: {final_metrics}")
        return final_metrics
    
    def _detect_repo_components(self) -> int:
        """Detect frontend components from repository structure"""
        component_count = 0
        
        # Check for component files in repository
        folder_structure = self._repo_analysis.get('folder_structure', {})
        logger.info(f"🔍 Analyzing folder structure: {list(folder_structure.keys())}")
        
        for folder, info in folder_structure.items():
            if isinstance(info, dict) and 'files' in info:
                folder_lower = folder.lower()
                files = info['files']
                logger.info(f"🔍 Folder '{folder}': {len(files)} files")
                
                # Count all frontend files, not just in specific folders
                for file in files:
                    if file.endswith(('.jsx', '.tsx', '.vue', '.js', '.ts', '.component.ts', '.html')):
                        component_count += 1
                        logger.info(f"✅ Found component file: {file}")
        
        # If no components found, check if we have any frontend files at all
        if component_count == 0:
            total_frontend_files = 0
            for folder, info in folder_structure.items():
                if isinstance(info, dict) and 'files' in info:
                    for file in info['files']:
                        if file.endswith(('.js', '.jsx', '.ts', '.tsx', '.vue', '.html', '.css')):
                            total_frontend_files += 1
            
            # Estimate components based on frontend files
            if total_frontend_files > 0:
                component_count = max(3, total_frontend_files // 2)  # Estimate
                logger.info(f"🔍 Estimated {component_count} components from {total_frontend_files} frontend files")
        
        logger.info(f"🎯 Total components detected: {component_count}")
        return component_count
    
    def _detect_repo_languages(self) -> List[str]:
        """Detect programming languages from repository files"""
        languages = set()
        
        # First check if languages are already provided in repo analysis
        existing_languages = self._repo_analysis.get('languages', {})
        if existing_languages:
            languages.update(existing_languages.keys())
            logger.info(f"🔍 Found existing languages: {list(languages)}")
        
        # Also scan folder structure for file extensions
        folder_structure = self._repo_analysis.get('folder_structure', {})
        for folder, info in folder_structure.items():
            if isinstance(info, dict) and 'files' in info:
                for file in info['files']:
                    if file.endswith('.py'): 
                        languages.add('Python')
                        logger.info(f"✅ Found Python file: {file}")
                    elif file.endswith(('.js', '.jsx')): 
                        languages.add('JavaScript')
                        logger.info(f"✅ Found JavaScript file: {file}")
                    elif file.endswith(('.ts', '.tsx')): 
                        languages.add('TypeScript')
                        logger.info(f"✅ Found TypeScript file: {file}")
                    elif file.endswith('.java'): 
                        languages.add('Java')
                        logger.info(f"✅ Found Java file: {file}")
                    elif file.endswith('.go'): 
                        languages.add('Go')
                        logger.info(f"✅ Found Go file: {file}")
                    elif file.endswith('.rs'): 
                        languages.add('Rust')
                        logger.info(f"✅ Found Rust file: {file}")
                    elif file.endswith('.php'): 
                        languages.add('PHP')
                        logger.info(f"✅ Found PHP file: {file}")
                    elif file.endswith('.rb'): 
                        languages.add('Ruby')
                        logger.info(f"✅ Found Ruby file: {file}")
                    elif file.endswith(('.html', '.css')): 
                        languages.add('HTML/CSS')
                        logger.info(f"✅ Found HTML/CSS file: {file}")
        
        result = list(languages)
        logger.info(f"🎯 Total languages detected: {result}")
        return result
    
    def _detect_repo_backend_tech(self) -> List[str]:
        """Detect backend technologies from repository"""
        backend_tech = []
        
        # Check for backend framework files
        for folder, info in self._repo_analysis.get('folder_structure', {}).items():
            if isinstance(info, dict) and 'files' in info:
                for file in info['files']:
                    if file in ['main.py', 'app.py', 'server.py']:
                        backend_tech.append('FastAPI/Flask')
                    elif file == 'package.json':
                        backend_tech.append('Node.js')
                    elif file in ['pom.xml', 'build.gradle']:
                        backend_tech.append('Spring Boot')
                    elif file == 'requirements.txt':
                        backend_tech.append('Python')
        
        return list(set(backend_tech))
    
    def _extract_prd_endpoints(self) -> List[Dict]:
        """Extract specific API endpoints mentioned in PRD"""
        endpoints = []
        prd_content = self._prd_analysis.get('content', '')
        
        if not prd_content:
            return endpoints
        
        # Find explicit endpoint definitions
        explicit_endpoints = re.findall(r'/(?:api/)?([a-zA-Z/]+)', prd_content)
        
        for endpoint_path in explicit_endpoints:
            if endpoint_path and len(endpoint_path) > 1:
                # Clean path
                clean_path = endpoint_path.strip('/')
                if '/' in clean_path:
                    parts = clean_path.split('/')
                    if parts[0] == 'auth':
                        method = 'POST' if parts[1] in ['register', 'login'] else 'GET'
                    else:
                        method = 'GET'
                else:
                    method = 'GET'
                
                endpoints.append({
                    'method': method,
                    'path': f'/api/{clean_path}' if not clean_path.startswith('auth') else f'/{clean_path}',
                    'purpose': f'{method} endpoint for {clean_path.replace("/", " ")}',
                    'auth_required': 'auth' not in clean_path,
                    'request_fields': self._get_endpoint_fields(clean_path, method, 'request'),
                    'response_fields': self._get_endpoint_fields(clean_path, method, 'response')
                })
        
        # Generate CRUD endpoints for database entities
        entities = self._extract_entities_from_prd(prd_content)
        for entity in entities:
            if entity not in [ep['path'].split('/')[-1] for ep in endpoints]:
                endpoints.extend(self._generate_entity_endpoints(entity))
        
        return endpoints
    
    def _extract_prd_backend_tech(self) -> List[str]:
        """Extract backend technologies from PRD or provide intelligent defaults"""
        backend_tech = []
        prd_content = self._prd_analysis.get('content', '').lower()
        
        # Check for explicit technology mentions
        tech_keywords = {
            'python': 'Python',
            'fastapi': 'FastAPI',
            'django': 'Django',
            'flask': 'Flask',
            'node': 'Node.js',
            'express': 'Express.js',
            'java': 'Java',
            'spring': 'Spring Boot',
            'go': 'Go',
            'rust': 'Rust'
        }
        
        for keyword, tech in tech_keywords.items():
            if keyword in prd_content:
                backend_tech.append(tech)
        
        # If no backend tech found but PRD has API/backend references, use Python defaults
        if not backend_tech and any(word in prd_content for word in ['api', 'backend', 'server', 'database']):
            backend_tech = ['Python', 'FastAPI', 'Pydantic']
        
        return backend_tech
    
    def _has_frontend_indicators(self) -> bool:
        """Check if repository has frontend indicators"""
        folder_structure = self._repo_analysis.get('folder_structure', {})
        
        # Check for frontend-related folders or files
        for folder, info in folder_structure.items():
            folder_lower = folder.lower()
            if any(pattern in folder_lower for pattern in ['src', 'app', 'frontend', 'client', 'ui']):
                return True
            
            if isinstance(info, dict) and 'files' in info:
                for file in info['files']:
                    if file in ['index.html', 'package.json', 'index.js', 'App.js', 'main.js']:
                        return True
        
        return False
    
    def _has_backend_indicators(self) -> bool:
        """Check if repository or PRD has backend indicators"""
        # Check repository
        folder_structure = self._repo_analysis.get('folder_structure', {})
        for folder, info in folder_structure.items():
            folder_lower = folder.lower()
            if any(pattern in folder_lower for pattern in ['api', 'server', 'backend', 'service']):
                return True
        
        # Check PRD content
        prd_content = self._prd_analysis.get('content', '').lower()
        return any(word in prd_content for word in ['api', 'backend', 'server', 'database', 'endpoint'])
    
    def _has_api_indicators(self) -> bool:
        """Check if there are API indicators in repository or PRD"""
        # Check PRD for API mentions
        prd_content = self._prd_analysis.get('content', '').lower()
        if any(word in prd_content for word in ['api', 'endpoint', 'rest', 'graphql']):
            return True
        
        # Check repository structure
        folder_structure = self._repo_analysis.get('folder_structure', {})
        for folder in folder_structure.keys():
            if any(pattern in folder.lower() for pattern in ['api', 'routes', 'controllers']):
                return True
        
        return False

    def _create_goals_scope(self) -> List:
        story = []
        story.append(Paragraph("2. Architecture Goals & Scope", self.styles['CustomHeading1']))
        
        # Architecture goals from PRD or intelligent defaults
        story.append(Paragraph("Architecture Goals:", self.styles['CustomHeading2']))
        goals = self._prd_analysis.get('goals', [])
        
        if not goals:
            # Generate context-aware goals
            if self._repo_analysis.get('frontend_tech'):
                goals.append("Deliver responsive and intuitive user interface")
            if self._repo_analysis.get('backend_tech'):
                goals.append("Implement scalable backend services")
            goals.extend([
                "Ensure secure and maintainable code architecture",
                "Support efficient development and deployment workflows",
                "Enable cross-platform compatibility and performance"
            ])
        
        for goal in goals[:5]:
            story.append(Paragraph(f"• {self._sanitize_text(goal)}", self.styles['CustomBullet']))
        
        return story

    def _create_context_diagram(self) -> List:
        story = []
        story.append(Paragraph("3. Architecture Context Diagram", self.styles['CustomHeading1']))
        
        # Context description
        context_text = "The following diagram illustrates the high-level context of the system, showing how users interact with the frontend, which communicates with backend services, and how data flows to external systems and databases."
        story.append(Paragraph(context_text, self.styles['CustomBody']))
        story.append(Spacer(1, 0.2*inch))
        
        try:
            # Generate visual context flow diagram
            diagram_generator = ArchitectureDiagramGenerator(self.output_dir)
            frontend_tech = self._repo_analysis.get('frontend_tech', ['React'])
            backend_tech = self._repo_analysis.get('backend_tech', ['FastAPI'])
            database_tech = self._repo_analysis.get('database_tech', ['PostgreSQL'])
            
            diagram_path = diagram_generator.generate_context_flow_diagram(
                frontend_tech=frontend_tech,
                backend_tech=backend_tech,
                database_tech=database_tech
            )
            
            # Embed diagram image
            from reportlab.platypus import Image
            if os.path.exists(diagram_path):
                img = Image(diagram_path, width=6.5*inch, height=4.5*inch)
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
        except Exception as e:
            logger.error(f"Failed to generate context diagram: {str(e)}")
            story.append(Paragraph("Context diagram generation failed.", self.styles['CustomBody']))
        
        # System boundaries
        story.append(Paragraph("System Boundaries:", self.styles['CustomHeading2']))
        fallback_boundaries = self._get_fallback_boundaries()
        for boundary in fallback_boundaries:
            story.append(Paragraph(f"• {boundary}", self.styles['CustomBullet']))
        
        return story

    def _create_frontend_section(self) -> List:
        story = []
        story.append(Paragraph("4. Frontend Architecture", self.styles['CustomHeading1']))
        
        # Framework info
        story.append(Paragraph("Framework & Structure:", self.styles['CustomHeading2']))
        frontend_tech = self._repo_analysis.get('frontend_tech', ['React'])
        story.append(Paragraph(f"Primary Framework: {', '.join(frontend_tech)}", self.styles['CustomBody']))
        
        # Component analysis
        story.append(Paragraph("Component Analysis:", self.styles['CustomHeading2']))
        components_count = self._repo_analysis.get('components_total', 0)
        story.append(Paragraph(f"Total Components: {components_count}", self.styles['CustomBody']))
        story.append(Paragraph(f"Total Pages: {self._repo_analysis.get('pages_total', 0)}", self.styles['CustomBody']))
        story.append(Spacer(1, 0.2*inch))
        
        try:
            # Generate visual frontend architecture diagram
            diagram_generator = ArchitectureDiagramGenerator(self.output_dir)
            component_names = self._extract_real_component_names()
            
            diagram_path = diagram_generator.generate_frontend_architecture_diagram(
                frontend_tech=frontend_tech,
                components_count=components_count,
                component_names=component_names
            )
            
            # Embed diagram image
            from reportlab.platypus import Image
            if os.path.exists(diagram_path):
                img = Image(diagram_path, width=6.5*inch, height=5.5*inch)
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
        except Exception as e:
            logger.error(f"Failed to generate frontend diagram: {str(e)}")
            story.append(Paragraph("Frontend diagram generation failed.", self.styles['CustomBody']))
        
        # Component structure analysis
        component_structure = self._analyze_frontend_structure()
        if component_structure:
            story.append(Paragraph("Component Structure:", self.styles['CustomHeading2']))
            for item in component_structure:
                story.append(Paragraph(f"• {item}", self.styles['CustomBullet']))
        
        return story

    def _create_backend_section(self) -> List:
        story = []
        
        has_backend = bool(self._repo_analysis.get('backend_tech'))
        title = "5. Backend Architecture" if has_backend else "5. Recommended Backend Architecture (No backend detected)"
        story.append(Paragraph(title, self.styles['CustomHeading1']))
        
        if not has_backend:
            story.append(Paragraph("No backend code detected in repository. Recommended architecture based on PRD:", self.styles['CustomBody']))
        
        # Framework info
        story.append(Paragraph("Framework & Structure:", self.styles['CustomHeading2']))
        backend_tech = self._repo_analysis.get('backend_tech', ['FastAPI (recommended)'])
        story.append(Paragraph(f"Primary Framework: {', '.join(backend_tech)}", self.styles['CustomBody']))
        
        # Get comprehensive API endpoints from all sources
        all_endpoints = self._get_comprehensive_api_endpoints()
        api_count = len(all_endpoints)
        
        story.append(Paragraph(f"API Endpoints ({api_count} detected):", self.styles['CustomHeading2']))
        
        if all_endpoints:
            for endpoint in all_endpoints[:15]:  # Show up to 15 endpoints
                if isinstance(endpoint, dict):
                    method = endpoint.get('method', 'GET')
                    path = endpoint.get('path', '/')
                    purpose = endpoint.get('purpose', '')
                    story.append(Paragraph(f"• {method} {path}", self.styles['CustomBullet']))
                    if purpose:
                        story.append(Paragraph(f"  Purpose: {purpose}", self.styles['CustomBody']))
                else:
                    story.append(Paragraph(f"• {endpoint}", self.styles['CustomBullet']))
        else:
            story.append(Paragraph("No API endpoints detected in repository or PRD", self.styles['CustomBody']))
        
        story.append(Spacer(1, 0.2*inch))
        
        try:
            # Generate visual backend architecture diagram
            diagram_generator = ArchitectureDiagramGenerator(self.output_dir)
            database_tech = self._repo_analysis.get('database_tech', ['PostgreSQL'])
            
            diagram_path = diagram_generator.generate_backend_architecture_diagram(
                backend_tech=backend_tech,
                api_count=api_count,
                database_tech=database_tech
            )
            
            # Embed diagram image
            from reportlab.platypus import Image
            if os.path.exists(diagram_path):
                img = Image(diagram_path, width=6.5*inch, height=5.5*inch)
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
        except Exception as e:
            logger.error(f"Failed to generate backend diagram: {str(e)}")
            story.append(Paragraph("Backend diagram generation failed.", self.styles['CustomBody']))
        
        # Service layer analysis
        service_analysis = self._analyze_backend_services(all_endpoints)
        if service_analysis:
            story.append(Paragraph("Service Layer Analysis:", self.styles['CustomHeading2']))
            for item in service_analysis:
                story.append(Paragraph(f"• {item}", self.styles['CustomBullet']))
        
        return story

    def _create_api_section(self) -> List:
        story = []
        story.append(Paragraph("6. API Endpoints Documentation", self.styles['CustomHeading1']))
        
        # Get existing endpoints
        repo_endpoints = self._repo_analysis.get('api_endpoints', [])
        prd_apis = self._prd_analysis.get('api_methods', [])
        
        # Generate intelligent endpoints based on frontend analysis
        inferred_endpoints = self._infer_api_endpoints()
        
        # Repository endpoints
        if repo_endpoints:
            story.append(Paragraph("Repository Detected Endpoints:", self.styles['CustomHeading2']))
            for endpoint in repo_endpoints:
                method = endpoint.get('method', 'GET')
                path = endpoint.get('path', '/')
                story.append(Paragraph(f"• {method} {path}", self.styles['CustomBullet']))
                story.append(Paragraph(f"  Description: {endpoint.get('description', 'No description')}", self.styles['CustomBody']))
        
        # PRD specified APIs
        if prd_apis:
            story.append(Paragraph("PRD Specified APIs:", self.styles['CustomHeading2']))
            for api in prd_apis:
                story.append(Paragraph(f"• {self._sanitize_text(api)}", self.styles['CustomBullet']))
        
        # Inferred endpoints with detailed specifications
        if inferred_endpoints:
            story.append(Paragraph("Inferred API Endpoints (Based on Frontend & PRD Analysis):", self.styles['CustomHeading2']))
            for endpoint in inferred_endpoints:
                # Endpoint header
                story.append(Paragraph(f"{endpoint['method']} {endpoint['path']}", self.styles['CustomHeading2']))
                story.append(Paragraph(f"Purpose: {endpoint['purpose']}", self.styles['CustomBody']))
                
                # Request fields table
                if endpoint.get('request_fields'):
                    story.append(Paragraph("Request Fields:", self.styles['CustomBullet']))
                    req_data = [['Field', 'Type', 'Required', 'Description']]
                    for field, details in endpoint['request_fields'].items():
                        req_data.append([
                            field,
                            details.get('type', 'string'),
                            'Yes' if details.get('required', False) else 'No',
                            details.get('description', 'No description')
                        ])
                    
                    req_table = Table(req_data, colWidths=[1.2*inch, 0.8*inch, 0.7*inch, 2.8*inch])
                    req_table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), self.colors['light_gray']),
                        ('TEXTCOLOR', (0, 0), (-1, 0), self.colors['text']),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                        ('GRID', (0, 0), (-1, -1), 1, self.colors['medium_gray']),
                    ]))
                    story.append(req_table)
                
                # Response fields table
                if endpoint.get('response_fields'):
                    story.append(Paragraph("Response Fields:", self.styles['CustomBullet']))
                    resp_data = [['Field', 'Type', 'Description']]
                    for field, details in endpoint['response_fields'].items():
                        resp_data.append([
                            field,
                            details.get('type', 'string'),
                            details.get('description', 'No description')
                        ])
                    
                    resp_table = Table(resp_data, colWidths=[1.5*inch, 1*inch, 3*inch])
                    resp_table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), self.colors['medium_gray']),
                        ('TEXTCOLOR', (0, 0), (-1, 0), self.colors['text']),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                        ('GRID', (0, 0), (-1, -1), 1, self.colors['medium_gray']),
                    ]))
                    story.append(resp_table)
                
                story.append(Spacer(1, 0.15*inch))
        
        # Summary
        total_endpoints = len(repo_endpoints) + len(prd_apis) + len(inferred_endpoints)
        story.append(Paragraph(f"Total API Endpoints: {total_endpoints}", self.styles['CustomHeading2']))
        
        return story

    def _create_data_flow_section(self) -> List:
        story = []
        
        # Get dynamic project title
        project_title = self._get_dynamic_project_title()
        story.append(Paragraph(f"7. {project_title} - Sequence Diagram", self.styles['CustomHeading1']))
        
        # Generate dynamic sequence diagram based on detected entities and APIs
        sequence_diagram = self._generate_dynamic_sequence_diagram()
        story.append(Paragraph("System Interaction Flow:", self.styles['CustomHeading2']))
        # Split diagram into lines and add each as separate paragraph to preserve formatting
        diagram_lines = sequence_diagram.strip().split('\n')
        for line in diagram_lines:
            story.append(Paragraph(line, self.styles['PlainASCII']))
        
        # Add flow analysis
        flow_analysis = self._analyze_system_flow()
        if flow_analysis:
            story.append(Paragraph("Flow Analysis:", self.styles['CustomHeading2']))
            for item in flow_analysis:
                story.append(Paragraph(f"• {item}", self.styles['CustomBullet']))
        
        return story
    
    def _get_dynamic_project_title(self) -> str:
        """Get dynamic project title from GitHub URL or PRD document"""
        # Try to get from PRD first
        prd_title = self._prd_analysis.get('product_name', '')
        if prd_title and prd_title != 'Application' and not prd_title.startswith('%PDF'):
            return prd_title
        
        # Fallback to generic title
        return "System Architecture"
    
    def _generate_dynamic_sequence_diagram(self) -> str:
        """Generate dynamic sequence diagram in professional ASCII art style with real data"""
        # Extract real data
        entities = self._extract_entities_from_prd(self._prd_analysis.get('content', ''))
        endpoints = self._get_comprehensive_api_endpoints()
        
        # Get system components with safe access
        frontend_tech_list = self._repo_analysis.get('frontend_tech', [])
        backend_tech_list = self._repo_analysis.get('backend_tech', [])
        database_tech_list = self._repo_analysis.get('database_tech', [])
        
        # Safe access with fallback defaults
        frontend_tech = frontend_tech_list[0] if frontend_tech_list else 'React'
        backend_tech = backend_tech_list[0] if backend_tech_list else 'Python FastAPI'
        database_tech = database_tech_list[0] if database_tech_list else 'PostgreSQL'
        
        # Determine main entity for flow
        main_entity = entities[0] if entities else 'data'
        
        # Get product name for title
        product_name = self._get_product_name()
        
        # Generate professional ASCII art sequence diagram
        diagram = f"""{product_name} - {main_entity.title()} Management Sequence Diagram (List + Create {main_entity.title()})

+-----------+     +-------------------+     +-------------------+     +-------------+
|   User    |     |    Frontend       |     |     Backend       |     |  Database   |
| (Browser) |     | ({frontend_tech:<15}) |     | ({backend_tech:<15}) |     | ({database_tech:<9})|
+-----------+     +-------------------+     +-------------------+     +-------------+
    |                   |                         |                         |
    | 1. Navigate to      |                         |                         |
    |    {main_entity} page        |                         |                         |
    |-------------------->|                         |                         |
    |                   |                         |                         |
    |                   | 2. GET /api/{main_entity}    |                         |
    |                   |------------------------>|                         |
    |                   |                         | 3. SELECT * FROM {main_entity}
    |                   |                         |------------------------>|
    |                   |                         |<------------------------|
    |                   |                         | 4. Return records       |
    |                   |<------------------------|                         |
    |                   | 5. 200 OK + JSON        |                         |
    |<--------------------|                         |                         |
    | 6. Display {main_entity} list |                         |                         |
    |                   |                         |                         |
    | 7. Click "Create"   |                         |                         |
    |-------------------->|                         |                         |
    |                   |                         |                         |
    |                   | 8. POST /api/{main_entity}   |                         |
    |                   | {{name, desc...}}        |                         |
    |                   |------------------------>|                         |
    |                   |                         | 9. INSERT INTO {main_entity}
    |                   |                         |------------------------>|
    |                   |                         |<------------------------|
    |                   |                         | 10. Return new ID       |
    |                   |<------------------------|                         |
    |                   | 11. 201 Created + JSON  |                         |
    |<--------------------|                         |                         |
    | 12. Success message |                         |                         |
    |    & update UI      |                         |                         |"""
        
        return diagram
    
    def _generate_entity_flow_steps(self, entity: str, endpoints: List[Dict], actors: List[str]) -> List[Dict]:
        """Generate flow steps based on main entity (hotel, booking, etc.)"""
        entity_endpoints = [ep for ep in endpoints if isinstance(ep, dict) and entity in ep.get('path', '').lower()]
        
        if not entity_endpoints:
            return self._generate_default_flow_steps(actors)
        
        steps = [
            {'from_idx': 0, 'to_idx': 1, 'message': f'Access {entity} management', 'description': f'User wants to manage {entity}'},
            {'from_idx': 1, 'to_idx': 2, 'message': f'GET /api/{entity}', 'description': f'Frontend fetches {entity} data'},
            {'from_idx': 2, 'to_idx': 3, 'message': f'Query {entity} records', 'description': f'Backend queries database'},
            {'from_idx': 3, 'to_idx': 2, 'message': f'Return {entity} data', 'description': 'Database returns results'},
            {'from_idx': 2, 'to_idx': 1, 'message': 'Return JSON', 'description': 'Backend sends data'},
            {'from_idx': 1, 'to_idx': 0, 'message': f'Display {entity} form', 'description': 'Frontend shows data'}
        ]
        
        # Add create/update flow if POST endpoint exists
        has_post = any(ep.get('method') == 'POST' for ep in entity_endpoints)
        if has_post:
            steps.extend([
                {'from_idx': 0, 'to_idx': 1, 'message': f'Submit {entity} form', 'description': f'User submits {entity} changes'},
                {'from_idx': 1, 'to_idx': 2, 'message': f'POST /api/{entity}', 'description': f'Frontend sends {entity} update'},
                {'from_idx': 2, 'to_idx': 3, 'message': f'Update {entity} record', 'description': f'Backend updates database'},
                {'from_idx': 3, 'to_idx': 2, 'message': 'Confirm update', 'description': 'Database confirms'},
                {'from_idx': 2, 'to_idx': 1, 'message': 'Return success', 'description': 'Backend confirms'},
                {'from_idx': 1, 'to_idx': 0, 'message': 'Show success', 'description': 'Frontend displays confirmation'}
            ])
        
        return steps
    
    def _generate_default_flow_steps(self, actors: List[str]) -> List[Dict]:
        """Generate default flow when no specific entities detected"""
        return [
            {'from_idx': 0, 'to_idx': 1, 'message': 'Request data form', 'description': 'User wants to manage data'},
            {'from_idx': 1, 'to_idx': 2, 'message': 'GET /api/data', 'description': 'Frontend fetches data'},
            {'from_idx': 2, 'to_idx': 3, 'message': 'Query records', 'description': 'Backend queries database'},
            {'from_idx': 3, 'to_idx': 2, 'message': 'Return records', 'description': 'Database returns results'},
            {'from_idx': 2, 'to_idx': 1, 'message': 'Return JSON', 'description': 'Backend sends data'},
            {'from_idx': 1, 'to_idx': 0, 'message': 'Display form', 'description': 'Frontend shows data'},
            {'from_idx': 0, 'to_idx': 1, 'message': 'Submit changes', 'description': 'User submits form'},
            {'from_idx': 1, 'to_idx': 2, 'message': 'POST /api/data', 'description': 'Frontend sends update'},
            {'from_idx': 2, 'to_idx': 3, 'message': 'Update record', 'description': 'Backend updates database'},
            {'from_idx': 3, 'to_idx': 2, 'message': 'Confirm update', 'description': 'Database confirms'},
            {'from_idx': 2, 'to_idx': 1, 'message': 'Return success', 'description': 'Backend confirms'},
            {'from_idx': 1, 'to_idx': 0, 'message': 'Show success', 'description': 'Frontend displays confirmation'}
        ]
    
    def _analyze_system_flow(self) -> List[str]:
        """Analyze system flow and return insights"""
        flow_insights = []
        
        # Analyze endpoints for flow patterns
        endpoints = self._get_comprehensive_api_endpoints()
        
        if not endpoints:
            flow_insights.append("No API endpoints detected - static frontend application")
            return flow_insights
        
        # Authentication flow analysis
        auth_endpoints = [ep for ep in endpoints if isinstance(ep, dict) and 
                         ('auth' in ep.get('path', '').lower() or 'login' in ep.get('path', '').lower())]
        if auth_endpoints:
            flow_insights.append(f"Authentication flow: {len(auth_endpoints)} auth-related endpoints")
        
        # Data flow analysis
        data_endpoints = [ep for ep in endpoints if isinstance(ep, dict) and 
                         ep.get('method') in ['GET', 'POST', 'PUT', 'DELETE']]
        if data_endpoints:
            methods = [ep.get('method') for ep in data_endpoints]
            flow_insights.append(f"Data flow operations: {', '.join(set(methods))}")
        
        # Frontend-backend communication
        frontend_tech = self._repo_analysis.get('frontend_tech', [])
        backend_tech = self._repo_analysis.get('backend_tech', [])
        if frontend_tech and backend_tech:
            flow_insights.append(f"Communication: {frontend_tech[0]} ↔ {backend_tech[0]} via REST API")
        
        return flow_insights

    def _create_interactions_section(self) -> List:
        story = []
        story.append(Paragraph("8. Component Interactions", self.styles['CustomHeading1']))
        
        story.append(Paragraph("Communication Patterns:", self.styles['CustomHeading2']))
        story.append(Paragraph("• RESTful API communication", self.styles['CustomBullet']))
        story.append(Paragraph("• JSON data exchange", self.styles['CustomBullet']))
        story.append(Paragraph("• HTTP/HTTPS protocols", self.styles['CustomBullet']))
        
        return story

    def _create_unified_diagram(self) -> List:
        story = []
        story.append(Paragraph("9. Database Details", self.styles['CustomHeading1']))
        
        # Architecture summary with bullet points
        story.append(Paragraph("Architecture Summary:", self.styles['CustomHeading2']))
        frontend = ', '.join(self._repo_analysis.get('frontend_tech', ['React', 'Vite', 'Next.js']))
        backend = ', '.join(self._repo_analysis.get('backend_tech', ['Python', 'FastAPI', 'Uvicorn']))
        database = ', '.join(self._repo_analysis.get('database_tech', ['PostgreSQL', 'SQLAlchemy']))
        build_tools = ', '.join(self._repo_analysis.get('build_tools', ['npm/yarn']))
        
        story.append(Paragraph(f"• Frontend: {frontend}", self.styles['CustomBullet']))
        story.append(Paragraph(f"• Backend: {backend}", self.styles['CustomBullet']))
        story.append(Paragraph(f"• Database: {database}", self.styles['CustomBullet']))
        if build_tools != 'npm/yarn':
            story.append(Paragraph(f"• Build Tools: {build_tools}", self.styles['CustomBullet']))
        
        # Enhanced database details
        story.append(Paragraph("Database Details:", self.styles['CustomHeading2']))
        db_details = self._generate_enhanced_database_details()
        for detail in db_details:
            story.append(Paragraph(f"• {detail}", self.styles['CustomBullet']))
        
        # Comprehensive database tables with full schema based on API endpoints
        story.append(Paragraph("Database Tables (Based on Detected API Endpoints):", self.styles['CustomHeading2']))
        table_schemas = self._generate_comprehensive_table_schemas()
        
        if table_schemas:
            for schema_line in table_schemas:
                if schema_line.startswith('  -') or schema_line.startswith('    -'):
                    # Field definitions and indexes
                    story.append(Paragraph(f"• {schema_line.strip()}", self.styles['CustomBullet']))
                elif schema_line.strip().endswith('Indexes:'):
                    # Index section header
                    story.append(Paragraph(f"• {schema_line.strip()}", self.styles['CustomBullet']))
                else:
                    # Table headers
                    story.append(Paragraph(f"• {schema_line}", self.styles['CustomBullet']))
        else:
            story.append(Paragraph("• No database tables detected from API endpoints", self.styles['CustomBullet']))
        
        return story
    
    def _generate_enhanced_database_details(self) -> List[str]:
        """Generate comprehensive database details with LLM analysis"""
        db_tech = self._repo_analysis.get('database_tech', ['PostgreSQL'])
        primary_db = db_tech[0] if db_tech else 'PostgreSQL'
        
        # Determine ORM based on backend technology
        backend_tech = self._repo_analysis.get('backend_tech', [])
        if 'Python' in str(backend_tech) or 'FastAPI' in str(backend_tech):
            orm = 'SQLAlchemy - Python SQL toolkit and Object-Relational Mapping'
        elif 'Node.js' in str(backend_tech):
            orm = 'Sequelize/Prisma - Node.js ORM for database operations'
        elif 'Java' in str(backend_tech):
            orm = 'Hibernate - Java persistence framework'
        else:
            orm = 'SQLAlchemy - Python SQL toolkit and Object-Relational Mapping'
        
        return [
            f"Primary Database: {primary_db} - Relational database for structured data storage",
            f"ORM Layer: {orm}",
            "Connection Pooling: Configured for optimal performance and resource management",
            "Transaction Management: ACID compliance with rollback capabilities",
            "Security: Encrypted connections, parameterized queries to prevent SQL injection",
            "Backup Strategy: Automated daily backups with point-in-time recovery",
            "Indexing: Optimized indexes on frequently queried columns for performance"
        ]
    
    def _generate_comprehensive_table_schemas(self) -> List[str]:
        """Generate comprehensive database table schemas based on detected API endpoints"""
        # Get all detected API endpoints
        all_endpoints = self._get_comprehensive_api_endpoints()
        
        # Extract entities from API endpoints
        entities_from_apis = self._extract_entities_from_endpoints(all_endpoints)
        
        # Combine with PRD tables
        prd_tables = self._prd_analysis.get('database_tables', [])
        
        # Merge and prioritize API-based entities
        all_entities = list(set(entities_from_apis + prd_tables))
        
        # If no entities found, use intelligent defaults
        if not all_entities:
            all_entities = self._infer_tables_from_prd_analysis()
        
        schema_lines = []
        for entity in all_entities[:6]:  # Limit to 6 most important tables
            table_schema = self._generate_api_based_table_schema(entity, all_endpoints)
            schema_lines.extend(table_schema)
        
        return schema_lines
    
    def _infer_tables_from_prd_analysis(self) -> List[str]:
        """Intelligently infer database tables from PRD content and features"""
        tables = []
        prd_content = self._prd_analysis.get('content', '').lower()
        prd_features = [f.lower() for f in self._prd_analysis.get('features', [])]
        
        # Analyze PRD content for entity patterns
        content_and_features = prd_content + ' ' + ' '.join(prd_features)
        
        # Smart entity detection with priority order
        entity_indicators = {
            'user': ['user', 'account', 'profile', 'authentication', 'login', 'member', 'customer'],
            'product': ['product', 'item', 'catalog', 'inventory', 'merchandise', 'goods'],
            'order': ['order', 'purchase', 'transaction', 'checkout', 'payment', 'booking'],
            'category': ['category', 'type', 'classification', 'group', 'section'],
            'session': ['session', 'login', 'auth', 'token', 'authentication'],
            'review': ['review', 'rating', 'feedback', 'comment', 'evaluation'],
            'notification': ['notification', 'alert', 'message', 'email', 'communication']
        }
        
        # Score entities based on frequency of indicators
        entity_scores = {}
        for entity, indicators in entity_indicators.items():
            score = sum(content_and_features.count(indicator) for indicator in indicators)
            if score > 0:
                entity_scores[entity] = score
        
        # Sort by score and take top entities
        sorted_entities = sorted(entity_scores.items(), key=lambda x: x[1], reverse=True)
        tables = [entity for entity, score in sorted_entities[:5]]
        
        # Ensure minimum essential tables
        if not tables:
            tables = ['user', 'session', 'data']
        elif 'user' not in tables:
            tables.insert(0, 'user')  # Always include user table
        
        return tables[:6]
    
    def _extract_entities_from_endpoints(self, endpoints: List[Dict]) -> List[str]:
        """Extract database entities from API endpoint paths"""
        entities = set()
        
        for endpoint in endpoints:
            if isinstance(endpoint, dict):
                path = endpoint.get('path', '')
                # Extract entity from API path
                path_parts = [p for p in path.split('/') if p and p != 'api']
                
                for part in path_parts:
                    # Clean entity name
                    clean_part = re.sub(r'[{}]', '', part).lower()
                    # Skip common non-entity parts
                    if clean_part not in ['auth', 'login', 'register', 'logout', 'health', 'status', 'id']:
                        if len(clean_part) > 2 and clean_part.isalpha():
                            # Convert to singular form for table name
                            singular = clean_part.rstrip('s') if clean_part.endswith('s') else clean_part
                            entities.add(singular)
        
        return list(entities)
    
    def _generate_api_based_table_schema(self, entity: str, endpoints: List[Dict]) -> List[str]:
        """Generate detailed table schema based on API endpoint analysis"""
        clean_name = str(entity).lower().strip()
        schema = []
        
        # Find related endpoints for this entity
        related_endpoints = self._find_entity_endpoints(clean_name, endpoints)
        
        # Extract fields from endpoint request/response data
        entity_fields = self._extract_fields_from_endpoints(clean_name, related_endpoints)
        
        # Generate table schema based on detected fields
        table_title = clean_name.replace('_', ' ').title()
        schema.append(f"{table_title} Table: {self._get_entity_description(clean_name)}")
        
        # Always include primary key
        schema.append("  - id (Primary Key, Auto-increment): Unique identifier")
        
        # Add detected fields
        for field_name, field_info in entity_fields.items():
            if field_name != 'id':  # Skip duplicate ID
                sql_type = self._convert_to_sql_type(field_info['type'])
                constraints = self._get_field_constraints(field_name, field_info)
                description = field_info.get('description', f'{table_title} {field_name}')
                schema.append(f"  - {field_name} ({sql_type}{constraints}): {description}")
        
        # Add standard audit fields
        schema.extend([
            "  - created_at (TIMESTAMP, NOT NULL, DEFAULT CURRENT_TIMESTAMP): Record creation time",
            "  - updated_at (TIMESTAMP, NULL, ON UPDATE CURRENT_TIMESTAMP): Last modification time"
        ])
        
        # Add indexes based on common query patterns
        indexes = self._generate_table_indexes(clean_name, entity_fields)
        if indexes:
            schema.append("  Indexes:")
            schema.extend([f"    - {idx}" for idx in indexes])
        
        return schema
    
    def _find_entity_endpoints(self, entity: str, endpoints: List[Dict]) -> List[Dict]:
        """Find all endpoints related to a specific entity"""
        related = []
        entity_variations = [entity, f"{entity}s", entity.rstrip('s')]
        
        for endpoint in endpoints:
            if isinstance(endpoint, dict):
                path = endpoint.get('path', '').lower()
                if any(var in path for var in entity_variations):
                    related.append(endpoint)
        
        return related
    
    def _extract_fields_from_endpoints(self, entity: str, endpoints: List[Dict]) -> Dict:
        """Extract database fields from endpoint request/response specifications"""
        fields = {}
        
        for endpoint in endpoints:
            if isinstance(endpoint, dict):
                # Extract from request fields (POST/PUT operations)
                request_fields = endpoint.get('request_fields', {})
                for field_name, field_info in request_fields.items():
                    if field_name not in ['page', 'limit', 'offset']:  # Skip pagination
                        fields[field_name] = field_info
                
                # Extract from response fields
                response_fields = endpoint.get('response_fields', {})
                for field_name, field_info in response_fields.items():
                    if field_name not in ['data', 'total', 'message', 'token']:  # Skip meta fields
                        fields[field_name] = field_info
        
        # Add entity-specific intelligent defaults if no fields detected
        if not fields:
            fields = self._get_intelligent_entity_fields(entity)
        
        return fields
    
    def _get_intelligent_entity_fields(self, entity: str) -> Dict:
        """Generate intelligent default fields based on entity type"""
        entity_lower = entity.lower()
        
        # Common fields for all entities
        fields = {
            'name': {'type': 'string', 'required': True, 'description': f'{entity.title()} name'},
            'description': {'type': 'string', 'required': False, 'description': f'{entity.title()} description'}
        }
        
        # Entity-specific fields based on common patterns
        if entity_lower in ['user', 'account', 'member']:
            fields.update({
                'email': {'type': 'string', 'required': True, 'description': 'User email address'},
                'password_hash': {'type': 'string', 'required': True, 'description': 'Encrypted password'},
                'first_name': {'type': 'string', 'required': False, 'description': 'User first name'},
                'last_name': {'type': 'string', 'required': False, 'description': 'User last name'},
                'is_active': {'type': 'boolean', 'required': True, 'description': 'Account status'}
            })
        elif entity_lower in ['product', 'item']:
            fields.update({
                'price': {'type': 'float', 'required': True, 'description': 'Product price'},
                'sku': {'type': 'string', 'required': False, 'description': 'Stock keeping unit'},
                'category_id': {'type': 'integer', 'required': False, 'description': 'Product category reference'},
                'stock_quantity': {'type': 'integer', 'required': True, 'description': 'Available stock'},
                'is_available': {'type': 'boolean', 'required': True, 'description': 'Product availability'}
            })
        elif entity_lower in ['order', 'purchase']:
            fields.update({
                'user_id': {'type': 'integer', 'required': True, 'description': 'Customer reference'},
                'total_amount': {'type': 'float', 'required': True, 'description': 'Order total value'},
                'status': {'type': 'string', 'required': True, 'description': 'Order status'},
                'payment_method': {'type': 'string', 'required': False, 'description': 'Payment method used'},
                'shipping_address': {'type': 'string', 'required': False, 'description': 'Delivery address'}
            })
        elif entity_lower in ['category', 'tag']:
            fields.update({
                'parent_id': {'type': 'integer', 'required': False, 'description': 'Parent category reference'},
                'slug': {'type': 'string', 'required': True, 'description': 'URL-friendly identifier'},
                'sort_order': {'type': 'integer', 'required': False, 'description': 'Display order'}
            })
        
        return fields
    
    def _convert_to_sql_type(self, field_type: str) -> str:
        """Convert API field type to SQL data type"""
        type_mapping = {
            'string': 'VARCHAR(255)',
            'integer': 'INT',
            'float': 'DECIMAL(10,2)',
            'boolean': 'BOOLEAN',
            'datetime': 'TIMESTAMP',
            'date': 'DATE',
            'text': 'TEXT',
            'object': 'JSON',
            'array': 'JSON'
        }
        return type_mapping.get(field_type.lower(), 'VARCHAR(255)')
    
    def _get_field_constraints(self, field_name: str, field_info: Dict) -> str:
        """Generate SQL constraints for field"""
        constraints = []
        
        if field_info.get('required', False):
            constraints.append('NOT NULL')
        
        # Special constraints based on field name
        field_lower = field_name.lower()
        if 'email' in field_lower:
            constraints.extend(['UNIQUE', 'NOT NULL'])
        elif field_lower in ['slug', 'sku']:
            constraints.append('UNIQUE')
        elif field_lower.endswith('_id'):
            constraints.append('FOREIGN KEY')
        
        return ', ' + ', '.join(constraints) if constraints else ''
    
    def _get_entity_description(self, entity: str) -> str:
        """Get descriptive text for entity table"""
        descriptions = {
            'user': 'User accounts and authentication data',
            'product': 'Product catalog and inventory management',
            'order': 'Order transactions and purchase records',
            'category': 'Product categorization and taxonomy',
            'review': 'User reviews and ratings',
            'session': 'User session management',
            'payment': 'Payment transactions and billing',
            'address': 'User addresses and shipping information',
            'cart': 'Shopping cart items and user selections'
        }
        return descriptions.get(entity.lower(), f'{entity.title()} data storage and management')
    
    def _generate_table_indexes(self, entity: str, fields: Dict) -> List[str]:
        """Generate recommended database indexes"""
        indexes = []
        
        # Common index patterns
        for field_name in fields.keys():
            field_lower = field_name.lower()
            
            # Foreign key indexes
            if field_lower.endswith('_id'):
                indexes.append(f"INDEX idx_{entity}_{field_name} ({field_name})")
            
            # Unique field indexes
            elif field_lower in ['email', 'slug', 'sku']:
                indexes.append(f"UNIQUE INDEX idx_{entity}_{field_name} ({field_name})")
            
            # Search field indexes
            elif field_lower in ['name', 'title']:
                indexes.append(f"INDEX idx_{entity}_{field_name} ({field_name})")
        
        # Entity-specific composite indexes
        if entity.lower() == 'user':
            indexes.append(f"INDEX idx_{entity}_active_email (is_active, email)")
        elif entity.lower() == 'product':
            indexes.append(f"INDEX idx_{entity}_category_price (category_id, price)")
        elif entity.lower() == 'order':
            indexes.append(f"INDEX idx_{entity}_user_status (user_id, status)")
        
        return indexes

    def _create_deployment_section(self) -> List:
        story = []
        story.append(Paragraph("10. Deployment Architecture", self.styles['CustomHeading1']))
        
        has_docker = 'Docker' in self._repo_analysis.get('build_tools', [])
        
        story.append(Paragraph("Deployment Configuration:", self.styles['CustomHeading2']))
        config = [
            f"Containerization: {'Yes (Docker detected)' if has_docker else 'Recommended'}",
            "Cloud Readiness: High",
            "Environment Configuration: Standard",
            "Monitoring Setup: Recommended"
        ]
        
        for item in config:
            story.append(Paragraph(f"• {item}", self.styles['CustomBullet']))
        
        return story

    def _create_security_section(self) -> List:
        story = []
        story.append(Paragraph("11. Security Model", self.styles['CustomHeading1']))
        
        story.append(Paragraph("Security Recommendations:", self.styles['CustomHeading2']))
        recommendations = [
            "Implement JWT token authentication",
            "Use HTTPS for all communications",
            "Validate and sanitize user inputs",
            "Implement rate limiting on APIs",
            "Regular security audits"
        ]
        
        for rec in recommendations:
            story.append(Paragraph(f"• {rec}", self.styles['CustomBullet']))
        
        return story

    def _create_tech_stack_section(self) -> List:
        story = []
        story.append(Paragraph("12. Technology Stack Summary", self.styles['CustomHeading1']))
        
        # Frontend technologies
        frontend = self._repo_analysis.get('frontend_tech', [])
        if frontend:
            story.append(Paragraph("Frontend Technologies:", self.styles['CustomHeading2']))
            for tech in frontend:
                story.append(Paragraph(f"• {tech}", self.styles['CustomBullet']))
        
        # Backend technologies
        backend = self._repo_analysis.get('backend_tech', [])
        if backend:
            story.append(Paragraph("Backend Technologies:", self.styles['CustomHeading2']))
            for tech in backend:
                story.append(Paragraph(f"• {tech}", self.styles['CustomBullet']))
        
        # Database technologies
        database = self._repo_analysis.get('database_tech', [])
        if database:
            story.append(Paragraph("Database Technologies:", self.styles['CustomHeading2']))
            for tech in database:
                story.append(Paragraph(f"• {tech}", self.styles['CustomBullet']))
        
        return story

    def _create_business_alignment(self) -> List:
        story = []
        story.append(Paragraph("13. Business Alignment", self.styles['CustomHeading1']))
        
        # PRD features
        features = self._prd_analysis.get('features', [])
        if features:
            story.append(Paragraph("Key Features (from PRD):", self.styles['CustomHeading2']))
            for feature in features:
                story.append(Paragraph(f"• {self._sanitize_text(feature)}", self.styles['CustomBullet']))
        
        return story

    def _create_recommendations_section(self) -> List:
        story = []
        story.append(Paragraph("14. Recommendations & Next Steps", self.styles['CustomHeading1']))
        
        story.append(Paragraph("Architecture Recommendations:", self.styles['CustomHeading2']))
        recommendations = self._get_recommendations()
        
        for i, rec in enumerate(recommendations, 1):
            story.append(Paragraph(f"{i}. {rec}", self.styles['CustomBody']))
        
        return story

    def _get_recommendations(self) -> List[str]:
        """Generate dynamic recommendations based on analysis"""
        recommendations = []
        
        # Based on missing technologies
        if not self._repo_analysis.get('backend_tech'):
            recommendations.append("Implement backend API service for data management")
        
        if not self._repo_analysis.get('database_tech'):
            recommendations.append("Add database layer for persistent data storage")
        
        if 'Docker' not in self._repo_analysis.get('build_tools', []):
            recommendations.append("Implement containerization with Docker")
        
        if not any('test' in tool.lower() for tool in self._repo_analysis.get('build_tools', [])):
            recommendations.append("Add comprehensive testing framework")
        
        # General recommendations
        recommendations.extend([
            "Implement monitoring and logging",
            "Add error handling and validation",
            "Consider implementing caching strategies",
            "Regular security audits and updates"
        ])
        
        return recommendations[:8]  # Limit to 8 recommendations
    
    def _create_architecture_diagram_section(self) -> List:
        """Create system architecture diagram section"""
        story = []
        story.append(Paragraph("15. System Architecture Diagram", self.styles['CustomHeading1']))
        
        try:
            # Generate architecture diagram
            diagram_generator = ArchitectureDiagramGenerator(self.output_dir)
            endpoints = self._get_comprehensive_api_endpoints()
            
            diagram_path = diagram_generator.generate_system_architecture_diagram(
                self._repo_analysis, self._prd_analysis, endpoints
            )
            
            # Add diagram description
            story.append(Paragraph("System Architecture Overview:", self.styles['CustomHeading2']))
            story.append(Paragraph("The following diagram shows the complete system architecture with real components, API endpoints, and data flow based on the analyzed repository and PRD.", self.styles['CustomBody']))
            story.append(Spacer(1, 0.2*inch))
            
            # Add diagram to PDF
            from reportlab.platypus import Image
            if os.path.exists(diagram_path):
                img = Image(diagram_path, width=6*inch, height=4.3*inch)
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
            
            # Add diagram insights
            story.append(Paragraph("Architecture Insights:", self.styles['CustomHeading2']))
            insights = self._generate_architecture_insights(endpoints)
            for insight in insights:
                story.append(Paragraph(f"• {insight}", self.styles['CustomBullet']))
                
        except Exception as e:
            logger.error(f"Failed to generate architecture diagram: {str(e)}")
            story.append(Paragraph("Architecture diagram generation failed. Please ensure matplotlib is installed.", self.styles['CustomBody']))
        
        return story
    
    def _create_layered_dataflow_section(self) -> List:
        """Create layered data flow diagram section (comprehensive multi-layer architecture)"""
        story = []
        story.append(Paragraph("16. Layered Data Flow Architecture", self.styles['CustomHeading1']))
        
        try:
            # Generate layered data flow diagram
            layered_generator = LayeredDataFlowGenerator(self.output_dir)
            endpoints = self._get_comprehensive_api_endpoints()
            
            diagram_path = layered_generator.generate_layered_dataflow_diagram(
                self._repo_analysis, self._prd_analysis, endpoints
            )
            
            # Add description
            product_name = self._prd_analysis.get('product_name', 'System')
            story.append(Paragraph("Comprehensive Architecture Layers:", self.styles['CustomHeading2']))
            story.append(Paragraph(
                f"This diagram illustrates the complete layered architecture of {product_name}, showing how data flows "
                f"through different architectural layers from data acquisition to presentation. Each layer is dynamically "
                f"detected from the repository analysis and PRD document.",
                self.styles['CustomBody']
            ))
            story.append(Spacer(1, 0.2*inch))
            
            # Add diagram to PDF
            from reportlab.platypus import Image
            if os.path.exists(diagram_path):
                img = Image(diagram_path, width=7*inch, height=5*inch)
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
            
            # Add layer descriptions
            story.append(Paragraph("Architecture Layers Explained:", self.styles['CustomHeading2']))
            
            layer_descriptions = [
                "<b>Presentation Layer:</b> User interface components, dashboards, and visualization elements that users interact with directly.",
                "<b>Application Layer:</b> Business logic, REST APIs, and core application services that process requests and orchestrate operations.",
                "<b>Knowledge Layer:</b> Analytics, machine learning models, and data processing engines that derive insights from data.",
                "<b>Acquisition Layer:</b> Data collection mechanisms including user inputs, IoT sensors, batch processes, and external data sources.",
                "<b>Storage Layer:</b> Persistent data storage including databases, caches, and file storage systems.",
                "<b>Security Layer:</b> Authentication, authorization, encryption, and security controls protecting the system.",
                "<b>CI/CD Layer:</b> Continuous integration and deployment tools for building, testing, and deploying the application.",
                "<b>External Systems:</b> Third-party APIs and external services integrated with the application."
            ]
            
            for desc in layer_descriptions:
                story.append(Paragraph(f"• {desc}", self.styles['CustomBullet']))
                story.append(Spacer(1, 0.05*inch))
                
        except Exception as e:
            logger.error(f"Failed to generate layered dataflow diagram: {str(e)}")
            story.append(Paragraph(f"Layered dataflow diagram generation failed: {str(e)}", self.styles['CustomBody']))
        
        return story

    def _create_sequence_diagram_section(self) -> List:
        """Create the new Section 7: System Interaction Flow (Sequence Diagram)"""
        story = []
        story.append(Paragraph("7. System Interaction Flow (Sequence Diagram)", self.styles['CustomHeading1']))
        
        try:
            # 1. DYNAMIC ENTITY DETECTION - Extract from multiple sources
            selected_entity = None
            api_route = None
            
            # Source 1: Extract from PRD Features (highest priority)
            prd_features = self._prd_analysis.get('features', [])
            prd_content = self._prd_analysis.get('content', '').lower()
            product_name = self._prd_analysis.get('product_name', 'Application')
            
            # Analyze PRD content for key nouns (entities)
            entity_candidates = set()
            
            # Extract from product name (e.g., "World Clock" -> "Clock")
            if product_name and product_name != 'Application':
                name_parts = product_name.split()
                for part in name_parts:
                    clean_part = part.strip().capitalize()
                    if len(clean_part) > 3 and clean_part.lower() not in ['system', 'application', 'app', 'platform']:
                        entity_candidates.add(clean_part)
            
            # Extract from features (e.g., "Manage timezones" -> "Timezone")
            for feature in prd_features[:5]:  # Top 5 features
                if isinstance(feature, str):
                    # Look for patterns like "manage X", "create X", "view X", "X management"
                    feature_lower = feature.lower()
                    patterns = [
                        r'manage\s+(\w+)',
                        r'create\s+(\w+)',
                        r'view\s+(\w+)',
                        r'add\s+(\w+)',
                        r'(\w+)\s+management',
                        r'(\w+)\s+list',
                        r'display\s+(\w+)',
                    ]
                    import re
                    for pattern in patterns:
                        matches = re.findall(pattern, feature_lower)
                        for match in matches:
                            clean_match = match.strip().capitalize()
                            if len(clean_match) > 3:
                                entity_candidates.add(clean_match)
            
            # Source 2: Extract from Database Tables
            db_tables = self._prd_analysis.get('database_tables', [])
            for table in db_tables[:3]:  # Top 3 tables
                if isinstance(table, str):
                    clean_table = table.strip().capitalize()
                    if len(clean_table) > 3:
                        entity_candidates.add(clean_table)
            
            # Source 3: Extract from API Endpoints
            endpoints = self._get_comprehensive_api_endpoints()
            endpoint_entities = set()
            for ep in endpoints:
                if isinstance(ep, dict):
                    path = ep.get('path', '')
                    # Extract entity from path like /api/clocks, /api/timezones
                    parts = [p for p in path.split('/') if p and p not in ['api', 'v1', 'v2', 'auth', 'login', 'register']]
                    if parts:
                        entity = parts[0].capitalize()
                        if len(entity) > 3:
                            endpoint_entities.add(entity)
            
            entity_candidates.update(endpoint_entities)
            
            # 4. INTELLIGENT SELECTION - Pick most relevant entity
            # Remove common generic terms
            generic_terms = {'data', 'item', 'items', 'object', 'record', 'entry', 'info', 'information'}
            entity_candidates = {e for e in entity_candidates if e.lower() not in generic_terms}
            
            if entity_candidates:
                # Prefer entities that appear in both PRD and endpoints
                prd_entities = {e for e in entity_candidates if any(e.lower() in f.lower() for f in prd_features) or e.lower() in prd_content}
                
                if prd_entities and endpoint_entities:
                    # Best case: entity in both PRD and endpoints
                    common = prd_entities.intersection(endpoint_entities)
                    if common:
                        selected_entity = list(common)[0]
                    else:
                        # Use PRD entity if available
                        selected_entity = list(prd_entities)[0]
                elif endpoint_entities:
                    selected_entity = list(endpoint_entities)[0]
                else:
                    selected_entity = list(entity_candidates)[0]
            
            # Fallback: Use product name or generic
            if not selected_entity:
                if product_name and product_name != 'Application':
                    # Extract last meaningful word from product name
                    words = product_name.split()
                    selected_entity = words[-1] if words else 'Resource'
                else:
                    selected_entity = 'Resource'
            
            # 5. Find matching API route
            selected_entity_lower = selected_entity.lower()
            api_route = f"/api/{selected_entity_lower}"
            
            # Try to find actual endpoint that matches
            for ep in endpoints:
                if isinstance(ep, dict):
                    path = ep.get('path', '')
                    if selected_entity_lower in path.lower():
                        api_route = path
                        break
            
            # 6. Generate Dynamic Diagram based on actual endpoints
            diagram_generator = ArchitectureDiagramGenerator(self.output_dir)
            
            diagram_path = diagram_generator.generate_sequence_diagram(
                project_name=product_name,
                entity_name=selected_entity,
                endpoints=endpoints  # Pass actual endpoints for dynamic flow generation
            )
            
            # 7. Add context-aware description
            story.append(Paragraph(f"Interaction Scenario: {selected_entity} Management (View & Create)", self.styles['CustomHeading2']))
            story.append(Paragraph(f"This sequence diagram illustrates the end-to-end flow for managing {selected_entity} within the {product_name}. It demonstrates the interaction between the User, Frontend, Backend, and Database layers.", self.styles['CustomBody']))
            story.append(Spacer(1, 0.1*inch))
            
            # 8. Embed image
            from reportlab.platypus import Image
            if os.path.exists(diagram_path):
                img = Image(diagram_path, width=7*inch, height=5*inch)
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
            
            # 9. Step-by-step text explanation
            story.append(Paragraph("Flow Description:", self.styles['CustomHeading2']))
            steps = [
                f"1. <b>User Action:</b> User navigates to the {selected_entity.lower()} page in the browser.",
                f"2. <b>Frontend Fetch:</b> The frontend initiates a GET request to <code>{api_route}</code>.",
                f"3. <b>Database Query:</b> The backend executes a <code>SELECT * FROM {selected_entity.lower()}</code> query.",
                f"4. <b>Data Render:</b> The list of records is returned and displayed to the user.",
                f"5. <b>User Action:</b> User clicks 'Create' and submits the form.",
                f"6. <b>Creation Request:</b> Frontend sends a POST request to <code>{api_route}</code> with payload.",
                f"7. <b>Persistence:</b> Backend executes <code>INSERT INTO {selected_entity.lower()}</code> to save the new record.",
                f"8. <b>Confirmation:</b> A success response (201 Created) is returned, and the UI updates."
            ]
            for step in steps:
                story.append(Paragraph(step, self.styles['CustomBody']))

        except Exception as e:
            logger.error(f"Failed to generate sequence diagram section: {str(e)}")
            story.append(Paragraph(f"Could not generate sequence diagram: {str(e)}", self.styles['CustomBody']))

        return story
    
    def _generate_architecture_insights(self, endpoints: List[Dict]) -> List[str]:
        """Generate insights about the architecture"""
        insights = []
        
        # Analyze endpoints
        if endpoints:
            entities = set()
            for ep in endpoints:
                if isinstance(ep, dict):
                    path = ep.get('path', '')
                    parts = [p for p in path.split('/') if p and p != 'api']
                    entities.update(parts)
            
            insights.append(f"System manages {len(entities)} main entities through {len(endpoints)} API endpoints")
        
        # Tech stack insights
        frontend = self._repo_analysis.get('frontend_tech', [])
        backend = self._repo_analysis.get('backend_tech', [])
        database = self._repo_analysis.get('database_tech', [])
        
        if frontend and backend:
            insights.append(f"Full-stack architecture: {frontend[0]} frontend communicates with {backend[0]} backend")
        
        if database:
            insights.append(f"Data persistence handled by {database[0]} with ORM integration")
        
        # Architecture pattern
        if len(endpoints) > 10:
            insights.append("RESTful API-driven architecture with clear separation of concerns")
        
        return insights[:5]
    
    def _generate_database_details(self) -> List[str]:
        """Generate detailed database information based on PRD analysis"""
        details = []
        
        # Get database technologies
        database_tech = self._repo_analysis.get('database_tech', ['PostgreSQL', 'SQLAlchemy'])
        
        # Primary database information
        primary_db = database_tech[0] if database_tech else 'PostgreSQL'
        details.append(f"Primary Database: {primary_db} - Relational database for structured data storage")
        
        # ORM/Database access layer
        if 'SQLAlchemy' in database_tech or not database_tech:
            details.append("ORM Layer: SQLAlchemy - Python SQL toolkit and Object-Relational Mapping")
        
        # Connection and performance details
        details.append("Connection Pooling: Configured for optimal performance and resource management")
        details.append("Transaction Management: ACID compliance with rollback capabilities")
        
        # Security features
        details.append("Security: Encrypted connections, parameterized queries to prevent SQL injection")
        
        # Backup and maintenance
        details.append("Backup Strategy: Automated daily backups with point-in-time recovery")
        details.append("Indexing: Optimized indexes on frequently queried columns for performance")
        
        return details
    
    def _generate_database_tables_info(self) -> List[str]:
        """Generate database tables information based on API endpoint analysis"""
        # Use the new comprehensive table schema generation
        return self._generate_comprehensive_table_schemas()
    
    def _infer_database_tables_from_prd(self) -> List[str]:
        """Infer database tables from PRD content using LLM analysis"""
        tables = []
        prd_content = self._prd_analysis.get('content', '').lower()
        
        # Common entity patterns to look for
        entity_patterns = {
            'users': ['user', 'account', 'profile', 'member', 'customer'],
            'products': ['product', 'item', 'catalog', 'inventory'],
            'orders': ['order', 'purchase', 'transaction', 'sale'],
            'categories': ['category', 'type', 'classification'],
            'sessions': ['session', 'login', 'authentication'],
            'settings': ['setting', 'configuration', 'preference'],
            'logs': ['log', 'audit', 'history', 'activity'],
            'notifications': ['notification', 'alert', 'message']
        }
        
        # Check PRD content for entity mentions
        for table_name, keywords in entity_patterns.items():
            if any(keyword in prd_content for keyword in keywords):
                tables.append(table_name)
        
        # Ensure we have at least basic tables
        if not tables:
            tables = ['users', 'data', 'sessions']
        
        return tables
    
    def _generate_table_schema(self, table_name: str) -> List[str]:
        """Generate detailed schema information for a database table"""
        schema_info = []
        
        # Clean table name
        clean_name = str(table_name).lower().strip()
        
        # Generate schema based on table type
        if clean_name in ['users', 'user']:
            schema_info.append("Users Table: User account and profile information")
            schema_info.append("  - id (Primary Key): Unique user identifier")
            schema_info.append("  - email (Unique): User email address for authentication")
            schema_info.append("  - password_hash: Encrypted password storage")
            schema_info.append("  - created_at: Account creation timestamp")
            schema_info.append("  - updated_at: Last profile update timestamp")
            
        elif clean_name in ['products', 'product', 'items', 'item']:
            schema_info.append("Products Table: Product catalog and inventory")
            schema_info.append("  - id (Primary Key): Unique product identifier")
            schema_info.append("  - name: Product name and title")
            schema_info.append("  - description: Detailed product description")
            schema_info.append("  - price: Product pricing information")
            schema_info.append("  - category_id (Foreign Key): Product category reference")
            
        elif clean_name in ['orders', 'order']:
            schema_info.append("Orders Table: Order and transaction records")
            schema_info.append("  - id (Primary Key): Unique order identifier")
            schema_info.append("  - user_id (Foreign Key): Customer reference")
            schema_info.append("  - total_amount: Order total value")
            schema_info.append("  - status: Order processing status")
            schema_info.append("  - created_at: Order placement timestamp")
            
        elif clean_name in ['sessions', 'session']:
            schema_info.append("Sessions Table: User session management")
            schema_info.append("  - id (Primary Key): Session identifier")
            schema_info.append("  - user_id (Foreign Key): User reference")
            schema_info.append("  - token: Authentication token")
            schema_info.append("  - expires_at: Session expiration time")
            
        else:
            # Generic table schema
            table_title = clean_name.replace('_', ' ').title()
            schema_info.append(table_title + " Table: Data storage")
            schema_info.append("  - id (Primary Key): Unique record identifier")
            schema_info.append("  - name: Record name or title")
            schema_info.append("  - data: JSON/Text field for flexible data storage")
            schema_info.append("  - created_at: Record creation timestamp")
            schema_info.append("  - updated_at: Last modification timestamp")
        
        return schema_info
    
    def _get_default_database_schema(self) -> List[str]:
        """Provide default database schema when no tables are detected"""
        return [
            "Users Table: User account and authentication data",
            "  - id (Primary Key): Unique user identifier",
            "  - email (Unique): User email for login",
            "  - password_hash: Encrypted password storage",
            "  - created_at: Account creation timestamp",
            "Data Table: Application data storage",
            "  - id (Primary Key): Unique record identifier",
            "  - user_id (Foreign Key): Owner reference",
            "  - content: JSON data field",
            "  - created_at: Record creation timestamp",
            "Sessions Table: User session management",
            "  - id (Primary Key): Session identifier",
            "  - user_id (Foreign Key): User reference",
            "  - token: Authentication token",
            "  - expires_at: Session expiration time"
        ]
    
    def _generate_fallback_diagram(self) -> str:
        """Generate fallback diagram when dynamic generation fails"""
        
        frontend = ', '.join(self._repo_analysis.get('frontend_tech', ['React']))
        backend = ', '.join(self._repo_analysis.get('backend_tech', ['FastAPI']))
        database = ', '.join(self._repo_analysis.get('database_tech', ['PostgreSQL']))
        
        return f"""
                   ╔═══════════════════════╗
                   ║       ■ CLIENT ■      ║
                   ║ (User Browser/Mobile) ║
                   ╚═══════════╤═══════════╝
                               ║
                               ║ ▼ (HTTP Request / API Call)
                               ▼
                   ╔═══════════╧═══════════╗
                   ║      ■ FRONTEND ■     ║
                   ║     ({frontend})      ║
                   ║  (Presentation Logic) ║
                   ╚═══════════╤═══════════╝
                               ║
                               ║ ▼ (Data Fetch / AJAX)
                               ▼
                   ╔═══════════╧═══════════╗
                   ║       ■ BACKEND ■     ║
                   ║     ({backend})       ║
                   ║   (Business Logic)    ║
                   ╚═══════════╤═══════════╝
                               ║
                               ║ ▼ (SQL / NoSQL Query)
                               ▼
                   ╔═══════════╧═══════════╗
                   ║      ■ DATABASE ■     ║
                   ║     ({database})      ║
                   ║   (Persistent Data)   ║
                   ╚═══════════════════════╝"""
    
    def _get_fallback_boundaries(self) -> List[str]:
        """Generate fallback system boundaries"""
        
        frontend_tech = ', '.join(self._repo_analysis.get('frontend_tech', [])) or 'React (Vite SPA)'
        backend_tech = ', '.join(self._repo_analysis.get('backend_tech', [])) or 'FastAPI (Python)'
        database_tech = ', '.join(self._repo_analysis.get('database_tech', [])) or 'PostgreSQL (Relational)'
        
        return [
            f"Frontend: {frontend_tech}",
            f"Backend: {backend_tech}",
            f"Database: {database_tech}",
            "External Integrations: Authentication, APIs, Cloud Services"
        ]
    
    def _infer_api_endpoints(self) -> List[Dict]:
        """Generate API endpoints from actual PRD content and frontend analysis"""
        endpoints = []
        
        # 1. Extract from PRD content (now returns structured objects)
        prd_endpoints = self._extract_prd_endpoints()
        endpoints.extend(prd_endpoints)
        
        # 2. Analyze frontend files for API calls
        frontend_endpoints = self._analyze_frontend_for_apis()
        endpoints.extend(frontend_endpoints)
        
        return endpoints
    
    def _get_endpoint_fields(self, endpoint_path: str, method: str, field_type: str) -> Dict:
        """Get appropriate fields for endpoint based on path and method"""
        fields = {}
        
        if 'auth' in endpoint_path:
            if 'register' in endpoint_path:
                if field_type == 'request':
                    fields = {
                        'name': {'type': 'string', 'required': True, 'description': 'User full name'},
                        'email': {'type': 'string', 'required': True, 'description': 'User email'},
                        'password': {'type': 'string', 'required': True, 'description': 'User password'}
                    }
                else:
                    fields = {
                        'user_id': {'type': 'integer', 'description': 'Created user ID'},
                        'message': {'type': 'string', 'description': 'Registration success message'}
                    }
            elif 'login' in endpoint_path:
                if field_type == 'request':
                    fields = {
                        'email': {'type': 'string', 'required': True, 'description': 'User email'},
                        'password': {'type': 'string', 'required': True, 'description': 'User password'}
                    }
                else:
                    fields = {
                        'token': {'type': 'string', 'description': 'JWT authentication token'},
                        'user': {'type': 'object', 'description': 'User information'}
                    }
        elif method == 'GET' and field_type == 'response':
            entity = endpoint_path.split('/')[-1]
            fields = {
                'data': {'type': 'array', 'description': f'Array of {entity} objects'},
                'total': {'type': 'integer', 'description': 'Total count'}
            }
        elif method == 'POST' and field_type == 'response':
            fields = {
                'id': {'type': 'integer', 'description': 'Created record ID'},
                'message': {'type': 'string', 'description': 'Success message'}
            }
        
        return fields
    
    def _generate_entity_endpoints(self, entity: str) -> List[Dict]:
        """Generate CRUD endpoints for discovered entity with intelligent field inference"""
        entity_clean = re.sub(r'[^\w]', '', entity.lower())
        if not entity_clean or len(entity_clean) < 3:
            return []
        
        # Infer fields based on entity type
        fields = self._infer_entity_fields(entity_clean)
        
        return [
            {
                'method': 'GET',
                'path': f'/api/{entity_clean}',
                'purpose': f'Retrieve {entity_clean} records',
                'auth_required': True,
                'request_fields': {
                    'page': {'type': 'integer', 'required': False, 'description': 'Page number'},
                    'limit': {'type': 'integer', 'required': False, 'description': 'Records per page'}
                },
                'response_fields': {
                    'data': {'type': 'array', 'description': f'Array of {entity_clean} objects'},
                    'total': {'type': 'integer', 'description': 'Total count'},
                    **{field: details for field, details in fields.items()}
                }
            },
            {
                'method': 'POST',
                'path': f'/api/{entity_clean}',
                'purpose': f'Create new {entity_clean}',
                'auth_required': True,
                'request_fields': fields,
                'response_fields': {
                    'id': {'type': 'integer', 'description': f'Created {entity_clean} ID'},
                    'message': {'type': 'string', 'description': 'Success message'}
                }
            },
            {
                'method': 'PUT',
                'path': f'/api/{entity_clean}/{{id}}',
                'purpose': f'Update {entity_clean}',
                'auth_required': True,
                'request_fields': fields,
                'response_fields': {
                    'message': {'type': 'string', 'description': 'Update confirmation'}
                }
            }
        ]
    
    def _infer_entity_fields(self, entity: str) -> Dict:
        """Intelligently infer fields for entity based on context and common patterns"""
        fields = {
            'name': {'type': 'string', 'required': True, 'description': f'{entity.title()} name'},
            'description': {'type': 'string', 'required': False, 'description': f'{entity.title()} description'}
        }
        
        # Add entity-specific fields based on semantic analysis
        entity_lower = entity.lower()
        
        # Analyze PRD content for entity-specific attributes
        prd_content = self._prd_analysis.get('content', '').lower()
        
        # Look for attributes mentioned near this entity
        entity_context = self._extract_entity_context(entity_lower, prd_content)
        
        for attr in entity_context:
            if attr not in fields:
                fields[attr] = {
                    'type': self._infer_field_type(attr),
                    'required': False,
                    'description': f'{entity.title()} {attr}'
                }
        
        return fields
    
    def _extract_entity_context(self, entity: str, content: str) -> List[str]:
        """Extract attributes mentioned in context of entity"""
        attributes = []
        
        # Find sentences containing the entity
        sentences = re.split(r'[.!?]', content)
        
        for sentence in sentences:
            if entity in sentence:
                # Extract potential attributes (nouns following common patterns)
                attr_patterns = [
                    rf'{entity}\s+(?:has|includes|contains|with)\s+([a-zA-Z]+)',
                    rf'([a-zA-Z]+)\s+(?:of|for)\s+{entity}',
                    rf'{entity}\s+([a-zA-Z]+)',
                ]
                
                for pattern in attr_patterns:
                    matches = re.findall(pattern, sentence)
                    for match in matches:
                        if len(match) > 2 and match not in ['the', 'and', 'with', 'for']:
                            attributes.append(match.lower())
        
        return list(set(attributes))[:5]  # Limit to most relevant
    
    def _infer_field_type(self, field_name: str) -> str:
        """Infer field type based on field name"""
        field_lower = field_name.lower()
        
        if any(word in field_lower for word in ['id', 'count', 'number', 'quantity', 'amount']):
            return 'integer'
        elif any(word in field_lower for word in ['price', 'cost', 'rate', 'percentage']):
            return 'float'
        elif any(word in field_lower for word in ['date', 'time', 'created', 'updated']):
            return 'datetime'
        elif any(word in field_lower for word in ['email', 'url', 'link']):
            return 'string'
        elif any(word in field_lower for word in ['active', 'enabled', 'verified']):
            return 'boolean'
        else:
            return 'string'
    
    def _infer_request_fields(self, path: str, method: str) -> Dict:
        """Infer request fields based on path and method with intelligent analysis"""
        fields = {}
        
        # Extract entity from path
        path_parts = path.strip('/').split('/')
        entity = path_parts[-1] if path_parts else 'resource'
        
        if method in ['POST', 'PUT', 'PATCH']:
            # Use entity-specific field inference
            if entity != 'resource':
                fields = self._infer_entity_fields(entity)
            else:
                fields['data'] = {'type': 'object', 'required': True, 'description': 'Request data'}
        
        # Special handling for auth endpoints
        if 'auth' in path.lower() or 'login' in path.lower():
            fields = {
                'email': {'type': 'string', 'required': True, 'description': 'User email'},
                'password': {'type': 'string', 'required': True, 'description': 'User password'}
            }
        
        return fields
    
    def _infer_response_fields(self, path: str, method: str) -> Dict:
        """Infer response fields based on path and method with intelligent analysis"""
        fields = {}
        
        # Extract entity from path
        path_parts = path.strip('/').split('/')
        entity = path_parts[-1] if path_parts else 'resource'
        
        if 'auth' in path.lower() or 'login' in path.lower():
            fields = {
                'token': {'type': 'string', 'description': 'JWT authentication token'},
                'user': {'type': 'object', 'description': 'User information'},
                'expires_in': {'type': 'integer', 'description': 'Token expiration time'}
            }
        elif method == 'GET':
            if entity != 'resource':
                entity_fields = self._infer_entity_fields(entity)
                fields = {
                    'data': {'type': 'array', 'description': f'Array of {entity} objects'},
                    **entity_fields
                }
            else:
                fields['data'] = {'type': 'object', 'description': 'Response data'}
        elif method == 'POST':
            fields = {
                'id': {'type': 'integer', 'description': f'Created {entity} ID'},
                'message': {'type': 'string', 'description': 'Success message'}
            }
        
        return fields
    
    def _get_default_analysis(self) -> Dict:
        """Return empty analysis when repo_analysis fails"""
        return {
            'languages': {},
            'frontend_tech': [],
            'backend_tech': [],
            'database_tech': [],
            'build_tools': [],
            'components_total': 0,
            'pages_total': 0,
            'api_endpoints': [],
            'patterns': [],
            'file_count': 0,
            'total_lines': 0
        }